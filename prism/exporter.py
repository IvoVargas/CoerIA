"""Exportação do programa, recursos e rastreabilidade de uma sessão CoerIA."""

from __future__ import annotations

import base64
import csv
import io
import json
import logging
import os
import re
import shutil
import subprocess
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile, TemporaryDirectory
from typing import Any

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt as DocxPt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.parts.image import Image as PptxImage
from pptx.util import Inches, Pt

from .branding import APP_FULL_NAME, APP_NAME
from .models import (
    RESOURCE_PRACTICAL,
    RESOURCE_PRESENTATION,
    RESOURCE_TEST,
    RESOURCE_WORKSHEET,
)
from .quality import attach_quality_report
from .relationships import derive_alignment_rows


LOGGER = logging.getLogger(__name__)
DOCUMENT_FORMAT_WORD = "word"
DOCUMENT_FORMAT_LATEX = "latex"
SUPPORTED_DOCUMENT_FORMATS = (
    DOCUMENT_FORMAT_WORD,
    DOCUMENT_FORMAT_LATEX,
)


def latex_pdf_compilation_enabled() -> bool:
    return str(os.getenv("COERIA_LATEX_PDF_ENABLED", "false")).strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }


def _latex_compiler_path() -> str:
    configured = str(os.getenv("COERIA_LATEX_COMPILER", "pdflatex")).strip()
    compiler = shutil.which(configured) if configured else None
    if not compiler:
        raise ValueError(
            "A compilação PDF de LaTeX está ativa, mas o compilador configurado "
            "não está disponível no servidor."
        )
    return compiler


def _latex_compile_timeout_seconds() -> int:
    try:
        configured = int(os.getenv("COERIA_LATEX_TIMEOUT_SECONDS", "45"))
    except ValueError:
        configured = 45
    return min(120, max(5, configured))


def normalize_document_formats(
    document_formats: list[str] | tuple[str, ...] | str | None,
) -> tuple[str, ...]:
    """Valida os formatos dos documentos textuais incluídos no pacote final."""

    if document_formats is None:
        requested = [DOCUMENT_FORMAT_WORD]
    elif isinstance(document_formats, str):
        requested = [document_formats]
    else:
        requested = list(document_formats)
    normalized = [str(item or "").strip().lower() for item in requested]
    invalid = [item for item in normalized if item not in SUPPORTED_DOCUMENT_FORMATS]
    if invalid:
        raise ValueError(
            "Formato documental inválido. Escolha Word, LaTeX ou ambos."
        )
    selected = tuple(
        item for item in SUPPORTED_DOCUMENT_FORMATS if item in normalized
    )
    if not selected:
        raise ValueError("Escolha pelo menos um formato documental para exportação.")
    return selected


def _safe_stem(value: str) -> str:
    clean = re.sub(r"[^0-9A-Za-zÀ-ÿ_-]+", "_", value.strip(), flags=re.UNICODE)
    return clean.strip("_")[:80] or "sessao_coeria"


def _add_document_title(document: Document, title: str, subtitle: str) -> None:
    document.add_heading(title, level=0)
    paragraph = document.add_paragraph(subtitle)
    paragraph.style = document.styles["Subtitle"]


def _set_document_defaults(document: Document) -> None:
    document.styles["Normal"].font.name = "Aptos"
    document.styles["Normal"].font.size = DocxPt(11)


PROGRAM_TABLE_WIDTH_DXA = 9960
PROGRAM_BLUE = "1F4D78"
PROGRAM_INK = "0B2545"
PROGRAM_HEADER_FILL = "E8EEF5"
PROGRAM_BORDER = "B8C4D1"


def _set_run_font(run: Any, name: str, size: float, *, bold: bool = False) -> None:
    run.font.name = name
    run.font.size = DocxPt(size)
    run.font.bold = bold
    fonts = run._element.get_or_add_rPr().get_or_add_rFonts()
    fonts.set(qn("w:ascii"), name)
    fonts.set(qn("w:hAnsi"), name)


def _configure_program_document(document: Document) -> None:
    """Aplica o preset compact_reference_guide com uma substituição A4 institucional."""

    section = document.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(1.7)
    section.right_margin = Cm(1.7)
    section.bottom_margin = Cm(1.7)
    section.left_margin = Cm(1.7)
    section.header_distance = Cm(1.0)
    section.footer_distance = Cm(1.0)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = DocxPt(10.5)
    normal.paragraph_format.space_after = DocxPt(6)
    normal.paragraph_format.line_spacing = 1.25

    style_tokens = {
        "Title": (24, PROGRAM_INK, 0, 8),
        "Subtitle": (11, "52677A", 0, 12),
        "Heading 1": (16, "2E74B5", 18, 10),
        "Heading 2": (13, "2E74B5", 14, 7),
        "Heading 3": (12, PROGRAM_BLUE, 10, 5),
    }
    for style_name, (size, color, before, after) in style_tokens.items():
        style = document.styles[style_name]
        style.font.name = "Arial"
        style.font.size = DocxPt(size)
        style.font.color.rgb = RGBColor.from_string(color)
        style.paragraph_format.space_before = DocxPt(before)
        style.paragraph_format.space_after = DocxPt(after)
        style.paragraph_format.keep_with_next = True

    title_ppr = document.styles["Title"]._element.get_or_add_pPr()
    title_border = title_ppr.find(qn("w:pBdr"))
    if title_border is not None:
        title_ppr.remove(title_border)

    for style_name in ("List Bullet", "List Number"):
        style = document.styles[style_name]
        style.font.name = "Arial"
        style.font.size = DocxPt(10.5)
        style.paragraph_format.left_indent = Cm(0.95)
        style.paragraph_format.first_line_indent = Cm(-0.48)
        style.paragraph_format.space_after = DocxPt(4)
        style.paragraph_format.line_spacing = 1.25

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = header.add_run("CoerIA  |  Programa da Unidade Curricular")
    _set_run_font(run, "Arial", 8.5)
    run.font.color.rgb = RGBColor.from_string("52677A")

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run(
        "Documento editável construído a partir dos artefactos aprovados pelo docente"
    )
    _set_run_font(run, "Arial", 8)
    run.font.color.rgb = RGBColor.from_string("687787")


def _cell_margins(cell: Any, *, top: int = 80, start: int = 120, bottom: int = 80, end: int = 120) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def _set_cell_fill(cell: Any, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = tc_pr.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        tc_pr.append(shading)
    shading.set(qn("w:fill"), fill)


def _set_table_geometry(table: Any, widths_dxa: list[int]) -> None:
    if sum(widths_dxa) != PROGRAM_TABLE_WIDTH_DXA:
        raise ValueError("As larguras da tabela do programa devem totalizar 9960 DXA.")
    table.autofit = False
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table_pr = table._tbl.tblPr
    table_width = table_pr.first_child_found_in("w:tblW")
    table_width.set(qn("w:w"), str(PROGRAM_TABLE_WIDTH_DXA))
    table_width.set(qn("w:type"), "dxa")
    table_indent = table_pr.first_child_found_in("w:tblInd")
    if table_indent is None:
        table_indent = OxmlElement("w:tblInd")
        table_pr.append(table_indent)
    table_indent.set(qn("w:w"), "120")
    table_indent.set(qn("w:type"), "dxa")

    grid = table._tbl.tblGrid
    for child in list(grid):
        grid.remove(child)
    for width in widths_dxa:
        grid_col = OxmlElement("w:gridCol")
        grid_col.set(qn("w:w"), str(width))
        grid.append(grid_col)

    for row in table.rows:
        for cell, width in zip(row.cells, widths_dxa):
            cell.width = Cm(width / 1440 * 2.54)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _cell_margins(cell)
            tc_width = cell._tc.get_or_add_tcPr().first_child_found_in("w:tcW")
            tc_width.set(qn("w:w"), str(width))
            tc_width.set(qn("w:type"), "dxa")


def _format_table(table: Any, headers: list[str], widths_dxa: list[int]) -> None:
    table.style = "Table Grid"
    table.rows[0]._tr.get_or_add_trPr().append(OxmlElement("w:tblHeader"))
    for index, header in enumerate(headers):
        cell = table.rows[0].cells[index]
        cell.text = header
        _set_cell_fill(cell, PROGRAM_HEADER_FILL)
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            paragraph.paragraph_format.space_after = DocxPt(0)
            for run in paragraph.runs:
                _set_run_font(run, "Arial", 9, bold=True)
                run.font.color.rgb = RGBColor.from_string(PROGRAM_INK)
    _set_table_geometry(table, widths_dxa)
    for row in table.rows[1:]:
        row_properties = row._tr.get_or_add_trPr()
        if row_properties.find(qn("w:cantSplit")) is None:
            row_properties.append(OxmlElement("w:cantSplit"))
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_after = DocxPt(2)
                paragraph.paragraph_format.line_spacing = 1.1
                for run in paragraph.runs:
                    _set_run_font(run, "Arial", 9)


def _display(value: Any, fallback: str = "A confirmar pelo docente") -> str:
    text = str(value or "").strip()
    return text or fallback


def _add_program_metadata(document: Document, course: dict[str, Any]) -> None:
    total_hours = float(course.get("contact_hours", 0) or 0) + float(
        course.get("autonomous_hours", 0) or 0
    )
    rows = [
        ("Unidade curricular", _display(course.get("unit_name"))),
        ("Curso ou programa", _display(course.get("program_name"))),
        ("Tipo de formação", _display(course.get("program_type"))),
        ("Ano curricular", _display(course.get("academic_year"))),
        ("Semestre", _display(course.get("semester"))),
        (
            "CNAEF",
            " — ".join(
                item
                for item in (
                    str(course.get("cnaef_code", "")).strip(),
                    str(course.get("cnaef_name", "")).strip(),
                )
                if item
            )
            or "A confirmar pelo docente",
        ),
        ("ECTS", _display(course.get("ects_credits") or "")),
        ("Horas de contacto", _display(course.get("contact_hours") or "")),
        ("Trabalho autónomo", _display(course.get("autonomous_hours") or "")),
        ("Carga de trabalho total", f"{total_hours:g} horas" if total_hours else "A confirmar pelo docente"),
        ("Taxonomia selecionada", _display(course.get("taxonomy_type"))),
    ]
    table = document.add_table(rows=1, cols=2)
    for label, value in rows:
        cells = table.add_row().cells
        cells[0].text = label
        cells[1].text = value
        _set_cell_fill(cells[0], "F2F4F7")
        for run in cells[0].paragraphs[0].runs:
            _set_run_font(run, "Arial", 9.5, bold=True)
    table._element.remove(table.rows[0]._element)
    _set_table_geometry(table, [2500, 7460])


def _bibliography_entries(value: Any) -> list[str]:
    return [
        item.strip(" -•\t")
        for item in re.split(r"[\r\n]+", str(value or ""))
        if item.strip(" -•\t")
    ]


def _validate_program_export_state(state: dict[str, Any]) -> None:
    required = (
        "learning_outcomes",
        "curriculum_analysis",
        "assessment_activities",
        "teaching_activities",
        "pedagogical_design",
    )
    missing = [key for key in required if not state.get(key)]
    if missing:
        raise ValueError(
            "Não é possível exportar o programa antes de aprovar toda a estrutura curricular: "
            + ", ".join(missing)
            + "."
        )


_LATEX_REPLACEMENTS = {
    "\\": r"\textbackslash{}",
    "{": r"\{",
    "}": r"\}",
    "$": r"\$",
    "&": r"\&",
    "%": r"\%",
    "#": r"\#",
    "_": r"\_",
    "~": r"\textasciitilde{}",
    "^": r"\textasciicircum{}",
    '"': r"\textquotedbl{}",
    "<": r"\textless{}",
    ">": r"\textgreater{}",
    "[": "{[}",
    "]": "{]}",
    "°": r"\textdegree{}",
    "≤": r"$\leq$",
    "≥": r"$\geq$",
    "→": r"$\rightarrow$",
    "×": r"$\times$",
    "±": r"$\pm$",
    "—": "---",
    "–": "--",
    "\u00a0": " ",
}


def _latex_escape(value: Any) -> str:
    """Escapa texto livre para que não altere a estrutura do documento LaTeX."""

    text = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
    escaped = "".join(_LATEX_REPLACEMENTS.get(character, character) for character in text)
    return escaped.replace("\n", r"\newline{}" + "\n")


def _latex_itemize(items: list[Any], *, ordered: bool = False) -> str:
    populated_items = [item for item in items if str(item or "").strip()]
    if not populated_items:
        return r"\emph{A confirmar pelo docente.}"
    environment = "enumerate" if ordered else "itemize"
    lines = [f"\\begin{{{environment}}}"]
    lines.extend(f"\\item{{}} {_latex_escape(item)}" for item in populated_items)
    lines.append(f"\\end{{{environment}}}")
    return "\n".join(lines)


def _latex_table(
    headers: list[str],
    rows: list[list[Any]],
    widths: list[float],
) -> str:
    if len(headers) != len(widths):
        raise ValueError("A tabela LaTeX necessita de uma largura por coluna.")
    column_spec = "@{}" + "".join(
        r">{\raggedright\arraybackslash}p{" + f"{width:.3f}" + r"\textwidth}"
        for width in widths
    ) + "@{}"
    header = " & ".join(
        rf"\textbf{{{_latex_escape(value)}}}" for value in headers
    ) + r" \\"
    body = [
        " & ".join(_latex_escape(value) for value in row) + r" \\"
        for row in rows
    ]
    return "\n".join(
        [
            r"{\footnotesize",
            rf"\begin{{longtable}}{{{column_spec}}}",
            r"\toprule",
            header,
            r"\midrule",
            r"\endfirsthead",
            r"\toprule",
            header,
            r"\midrule",
            r"\endhead",
            r"\bottomrule",
            r"\endfoot",
            *body,
            r"\end{longtable}",
            "}",
        ]
    )


def _latex_document(title: str, subtitle: str, body: list[str]) -> str:
    return "\n".join(
        [
            r"\documentclass[11pt,a4paper]{article}",
            r"\usepackage[T1]{fontenc}",
            r"\usepackage[utf8]{inputenc}",
            r"\usepackage[portuguese]{babel}",
            r"\usepackage[a4paper,margin=2cm]{geometry}",
            r"\usepackage{array}",
            r"\usepackage{booktabs}",
            r"\usepackage{longtable}",
            r"\usepackage{textcomp}",
            r"\usepackage[hidelinks]{hyperref}",
            r"\setlength{\parindent}{0pt}",
            r"\setlength{\parskip}{0.55em}",
            r"\setlength{\tabcolsep}{3pt}",
            r"\sloppy",
            rf"\title{{{_latex_escape(title)}}}",
            rf"\author{{{_latex_escape(subtitle)}}}",
            r"\date{}",
            r"\begin{document}",
            r"\maketitle",
            *body,
            r"\end{document}",
            "",
        ]
    )


def _write_latex_document(
    content: str,
    output_path: Path | str | None,
    *,
    prefix: str,
) -> str:
    if output_path is None:
        with NamedTemporaryFile(prefix=prefix, suffix=".tex", delete=False) as temp_file:
            destination = Path(temp_file.name)
    else:
        destination = Path(output_path)
    destination.write_text(content, encoding="utf-8")
    return str(destination)


def _latex_compiler_log_excerpt(
    output: bytes | str | None,
    working_directory: Path,
    *,
    maximum_characters: int = 12_000,
) -> str:
    """Prepara uma saída limitada do compilador para o diário técnico."""

    if isinstance(output, bytes):
        text = output.decode("utf-8", errors="replace")
    else:
        text = str(output or "")
    text = re.sub(r"\x1b\[[0-9;]*m", "", text)
    text = text.replace(str(working_directory), "<diretório-temporário>").strip()
    if not text:
        return "(sem saída do compilador)"
    if len(text) > maximum_characters:
        return "[… saída anterior omitida …]\n" + text[-maximum_characters:]
    return text


def compile_latex_pdf(source_path: Path | str) -> str | None:
    """Compila um `.tex` controlado pela aplicação, quando ativado no servidor."""

    if not latex_pdf_compilation_enabled():
        return None
    source = Path(source_path).resolve()
    if source.suffix.lower() != ".tex" or not source.is_file():
        raise ValueError("O ficheiro LaTeX a compilar não está disponível.")
    compiler = _latex_compiler_path()
    command = [
        compiler,
        "-interaction=nonstopmode",
        "-halt-on-error",
        "-file-line-error",
        "-no-shell-escape",
        "-output-directory",
        str(source.parent),
        source.name,
    ]
    environment = os.environ.copy()
    environment.update(
        {
            "openin_any": "p",
            "openout_any": "p",
        }
    )
    last_output: bytes | str | None = None
    for pass_number in (1, 2):
        try:
            result = subprocess.run(
                command,
                cwd=source.parent,
                env=environment,
                stdin=subprocess.DEVNULL,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                timeout=_latex_compile_timeout_seconds(),
                check=False,
            )
        except subprocess.TimeoutExpired as error:
            last_output = error.stdout
            LOGGER.error(
                "Timeout na passagem %s/2 da compilação LaTeX de %s. Saída do "
                "compilador:\n%s",
                pass_number,
                source.name,
                _latex_compiler_log_excerpt(last_output, source.parent),
            )
            raise ValueError(
                "A compilação PDF excedeu o tempo máximo permitido. "
                "O pacote não foi criado."
            ) from error
        last_output = result.stdout
        if result.returncode != 0:
            LOGGER.error(
                "Falha na passagem %s/2 da compilação LaTeX de %s "
                "(código %s). Saída do compilador:\n%s",
                pass_number,
                source.name,
                result.returncode,
                _latex_compiler_log_excerpt(last_output, source.parent),
            )
            raise ValueError(
                "Não foi possível compilar o documento LaTeX para PDF. "
                "O pacote não foi criado; escolha Word ou contacte o responsável "
                "técnico."
            )
    destination = source.with_suffix(".pdf")
    pdf_data = destination.read_bytes() if destination.is_file() else b""
    if (
        not pdf_data.startswith(b"%PDF-")
        or b"%%EOF" not in pdf_data[-2048:]
    ):
        LOGGER.error(
            "A compilação LaTeX de %s terminou sem produzir um PDF válido. "
            "Última saída do compilador:\n%s",
            source.name,
            _latex_compiler_log_excerpt(last_output, source.parent),
        )
        raise ValueError(
            "Não foi possível compilar o documento LaTeX para PDF. "
            "O pacote não foi criado; escolha Word ou contacte o responsável técnico."
        )
    LOGGER.info("PDF LaTeX compilado em duas passagens: %s", source.name)
    return str(destination)


def export_program_document(
    state: dict[str, Any], output_path: Path | str | None = None
) -> str:
    """Constrói o programa editável da UC a partir de artefactos já aprovados."""

    _validate_program_export_state(state)
    alignment_rows = derive_alignment_rows(state)

    course = state.get("course", {})
    analysis = state.get("curriculum_analysis", {})
    document = Document()
    _configure_program_document(document)
    title = document.add_paragraph(style="Title")
    title.add_run("Programa da Unidade Curricular")
    subtitle = document.add_paragraph(style="Subtitle")
    subtitle.add_run(_display(course.get("unit_name")))

    document.add_heading("1. Identificação e carga de trabalho", level=1)
    _add_program_metadata(document, course)

    document.add_heading("2. Objetivos gerais", level=1)
    general_aims = str(
        analysis.get("objectives") or course.get("general_aims", "")
    ).strip()
    document.add_paragraph(general_aims or "A confirmar pelo docente.")

    document.add_heading("3. Conteúdos programáticos", level=1)
    contents = analysis.get("contents", [])
    table = document.add_table(rows=1, cols=3)
    for item in contents:
        cells = table.add_row().cells
        cells[0].text = str(item.get("id", ""))
        cells[1].text = str(item.get("title", ""))
        cells[2].text = str(item.get("description", ""))
    _format_table(table, ["ID", "Conteúdo", "Descrição"], [800, 2700, 6460])

    document.add_heading("4. Resultados de aprendizagem", level=1)
    table = document.add_table(rows=1, cols=5)
    for outcome in state.get("learning_outcomes", []):
        cells = table.add_row().cells
        cells[0].text = str(outcome.get("id", ""))
        cells[1].text = str(outcome.get("statement", ""))
        cells[2].text = str(course.get("taxonomy_type", ""))
        cells[3].text = str(outcome.get("taxonomy_level", ""))
        cells[4].text = str(outcome.get("action_verb", ""))
    _format_table(
        table,
        ["ID", "Resultado de aprendizagem", "Taxonomia", "Nível", "Verbo"],
        [700, 4660, 1200, 2300, 1100],
    )

    document.add_heading("5. Atividades de ensino-aprendizagem", level=1)
    table = document.add_table(rows=1, cols=7)
    for activity in state.get("teaching_activities", []):
        cells = table.add_row().cells
        cells[0].text = str(activity.get("id", ""))
        cells[1].text = str(activity.get("learning_context", ""))
        cells[2].text = str(activity.get("activity", ""))
        cells[3].text = str(
            activity.get("practice") or activity.get("method", "")
        )
        cells[4].text = str(activity.get("support", ""))
        cells[5].text = str(activity.get("feedback_strategy", ""))
        cells[6].text = ", ".join(activity.get("outcome_ids", []))
    _format_table(
        table,
        [
            "ID",
            "Contexto",
            "Atividade",
            "Prática",
            "Acompanhamento",
            "Feedback",
            "Resultados",
        ],
        [500, 1250, 1900, 1550, 1900, 1500, 1360],
    )

    document.add_heading("6. Tarefas e critérios de avaliação", level=1)
    table = document.add_table(rows=1, cols=6)
    for assessment in state.get("assessment_activities", []):
        cells = table.add_row().cells
        cells[0].text = str(assessment.get("id", ""))
        cells[1].text = str(assessment.get("assessment_purpose", ""))
        cells[2].text = str(assessment.get("work_type", ""))
        cells[3].text = str(assessment.get("activity", ""))
        cells[4].text = str(assessment.get("criterion", ""))
        cells[5].text = ", ".join(assessment.get("outcome_ids", []))
    _format_table(
        table,
        ["ID", "Finalidade", "Modalidade", "Atividade", "Critério", "Resultados"],
        [650, 1200, 1400, 2800, 2510, 1400],
    )

    document.add_heading("7. Organização da sequência pedagógica", level=1)
    pedagogical_design = state.get("pedagogical_design", {})
    strategy = document.add_paragraph()
    strategy.add_run("Estratégia pedagógica: ").bold = True
    strategy.add_run(_display(pedagogical_design.get("strategy")))
    table = document.add_table(rows=1, cols=4)
    for item in pedagogical_design.get("sequence", []):
        cells = table.add_row().cells
        cells[0].text = str(item.get("outcome_id", ""))
        cells[1].text = str(item.get("focus", ""))
        cells[2].text = str(item.get("teaching_activity", ""))
        cells[3].text = str(item.get("assessment", ""))
    _format_table(
        table,
        ["Resultado", "Foco", "Atividade de ensino-aprendizagem", "Avaliação"],
        [700, 3000, 3200, 3060],
    )

    document.add_heading("8. Síntese automática do alinhamento", level=1)
    table = document.add_table(rows=1, cols=6)
    for row in alignment_rows:
        cells = table.add_row().cells
        cells[0].text = str(row.get("outcome_id", ""))
        cells[1].text = ", ".join(row.get("content_ids", []))
        cells[2].text = ", ".join(row.get("teaching_activity_ids", []))
        cells[3].text = ", ".join(row.get("assessment_ids", []))
        cells[4].text = str(row.get("status", ""))
        cells[5].text = str(row.get("rationale", ""))
    _format_table(
        table,
        ["RA", "Conteúdos", "Ensino-aprendizagem", "Avaliação", "Estado", "Fundamentação"],
        [650, 1250, 1750, 1350, 1150, 3810],
    )

    document.add_heading("9. Bibliografia", level=1)
    bibliography = _bibliography_entries(course.get("bibliography"))
    if bibliography:
        for entry in bibliography:
            document.add_paragraph(entry, style="List Bullet")
    else:
        paragraph = document.add_paragraph(
            "Bibliografia a fornecer ou validar pelo docente antes da utilização institucional."
        )
        paragraph.runs[0].italic = True
        paragraph.runs[0].font.color.rgb = RGBColor.from_string("7A5A00")

    if output_path is None:
        with NamedTemporaryFile(prefix="coeria_programa_", suffix=".docx", delete=False) as temp_file:
            destination = Path(temp_file.name)
    else:
        destination = Path(output_path)
    document.save(destination)
    return str(destination)


def export_program_latex(
    state: dict[str, Any], output_path: Path | str | None = None
) -> str:
    """Constrói em LaTeX o mesmo programa da UC disponibilizado em Word."""

    _validate_program_export_state(state)
    alignment_rows = derive_alignment_rows(state)
    course = state.get("course", {})
    analysis = state.get("curriculum_analysis", {})
    total_hours = float(course.get("contact_hours", 0) or 0) + float(
        course.get("autonomous_hours", 0) or 0
    )
    metadata_rows = [
        ["Unidade curricular", _display(course.get("unit_name"))],
        ["Curso ou programa", _display(course.get("program_name"))],
        ["Tipo de formação", _display(course.get("program_type"))],
        ["Ano curricular", _display(course.get("academic_year"))],
        ["Semestre", _display(course.get("semester"))],
        [
            "CNAEF",
            " --- ".join(
                item
                for item in (
                    str(course.get("cnaef_code", "")).strip(),
                    str(course.get("cnaef_name", "")).strip(),
                )
                if item
            )
            or "A confirmar pelo docente",
        ],
        ["ECTS", _display(course.get("ects_credits") or "")],
        ["Horas de contacto", _display(course.get("contact_hours") or "")],
        ["Trabalho autónomo", _display(course.get("autonomous_hours") or "")],
        [
            "Carga de trabalho total",
            f"{total_hours:g} horas" if total_hours else "A confirmar pelo docente",
        ],
        ["Taxonomia selecionada", _display(course.get("taxonomy_type"))],
    ]
    body = [
        r"\section{Identificação e carga de trabalho}",
        _latex_table(["Campo", "Valor"], metadata_rows, [0.24, 0.66]),
        r"\section{Objetivos gerais}",
    ]
    general_aims = str(
        analysis.get("objectives") or course.get("general_aims", "")
    ).strip()
    body.append(_latex_escape(general_aims or "A confirmar pelo docente."))
    body.extend(
        [
            r"\section{Conteúdos programáticos}",
            _latex_table(
                ["ID", "Conteúdo", "Descrição"],
                [
                    [item.get("id", ""), item.get("title", ""), item.get("description", "")]
                    for item in analysis.get("contents", [])
                ],
                [0.07, 0.22, 0.59],
            ),
            r"\section{Resultados de aprendizagem}",
            _latex_table(
                ["ID", "Resultado de aprendizagem", "Taxonomia", "Nível", "Verbo"],
                [
                    [
                        outcome.get("id", ""),
                        outcome.get("statement", ""),
                        course.get("taxonomy_type", ""),
                        outcome.get("taxonomy_level", ""),
                        outcome.get("action_verb", ""),
                    ]
                    for outcome in state.get("learning_outcomes", [])
                ],
                [0.06, 0.35, 0.11, 0.18, 0.12],
            ),
            r"\section{Atividades de ensino-aprendizagem}",
            _latex_table(
                [
                    "ID",
                    "Contexto",
                    "Atividade",
                    "Prática",
                    "Acompanhamento",
                    "Feedback",
                    "Resultados",
                ],
                [
                    [
                        activity.get("id", ""),
                        activity.get("learning_context", ""),
                        activity.get("activity", ""),
                        activity.get("practice") or activity.get("method", ""),
                        activity.get("support", ""),
                        activity.get("feedback_strategy", ""),
                        ", ".join(activity.get("outcome_ids", [])),
                    ]
                    for activity in state.get("teaching_activities", [])
                ],
                [0.045, 0.100, 0.170, 0.130, 0.160, 0.120, 0.095],
            ),
            r"\section{Tarefas e critérios de avaliação}",
            _latex_table(
                ["ID", "Finalidade", "Modalidade", "Atividade", "Critério", "Resultados"],
                [
                    [
                        assessment.get("id", ""),
                        assessment.get("assessment_purpose", ""),
                        assessment.get("work_type", ""),
                        assessment.get("activity", ""),
                        assessment.get("criterion", ""),
                        ", ".join(assessment.get("outcome_ids", [])),
                    ]
                    for assessment in state.get("assessment_activities", [])
                ],
                [0.05, 0.10, 0.12, 0.17, 0.20, 0.14],
            ),
            r"\section{Organização da sequência pedagógica}",
            r"\textbf{Estratégia pedagógica:} "
            + _latex_escape(
                _display(state.get("pedagogical_design", {}).get("strategy"))
            ),
            _latex_table(
                ["Resultado", "Foco", "Atividade de ensino-aprendizagem", "Avaliação"],
                [
                    [
                        item.get("outcome_id", ""),
                        item.get("focus", ""),
                        item.get("teaching_activity", ""),
                        item.get("assessment", ""),
                    ]
                    for item in state.get("pedagogical_design", {}).get("sequence", [])
                ],
                [0.06, 0.25, 0.27, 0.24],
            ),
            r"\clearpage",
            r"\section{Síntese automática do alinhamento}",
            _latex_table(
                ["RA", "Conteúdos", "Ensino-aprendizagem", "Avaliação", "Estado", "Fundamentação"],
                [
                    [
                        row.get("outcome_id", ""),
                        ", ".join(row.get("content_ids", [])),
                        ", ".join(row.get("teaching_activity_ids", [])),
                        ", ".join(row.get("assessment_ids", [])),
                        row.get("status", ""),
                        row.get("rationale", ""),
                    ]
                    for row in alignment_rows
                ],
                [0.05, 0.12, 0.17, 0.12, 0.10, 0.29],
            ),
            r"\section{Bibliografia}",
        ]
    )
    bibliography = _bibliography_entries(course.get("bibliography"))
    if bibliography:
        body.append(_latex_itemize(bibliography))
    else:
        body.append(
            r"\emph{Bibliografia a fornecer ou validar pelo docente antes da utilização institucional.}"
        )
    content = _latex_document(
        "Programa da Unidade Curricular",
        _display(course.get("unit_name")),
        body,
    )
    return _write_latex_document(
        content,
        output_path,
        prefix="coeria_programa_",
    )


PPT_NAVY = PptxRGBColor(11, 37, 69)
PPT_BLUE = PptxRGBColor(46, 116, 181)
PPT_TEAL = PptxRGBColor(40, 157, 143)
PPT_GOLD = PptxRGBColor(232, 174, 49)
PPT_PALE = PptxRGBColor(237, 243, 248)
PPT_INK = PptxRGBColor(28, 45, 61)
PPT_WHITE = PptxRGBColor(255, 255, 255)
PPT_MUTED = PptxRGBColor(91, 108, 122)


def _set_shape_alt_text(shape: Any, description: str) -> None:
    """Guarda uma descrição acessível no XML não-visual da forma."""

    if not description:
        return
    properties = shape._element.xpath(".//p:cNvPr")
    if properties:
        properties[0].set("descr", description)


def _add_ppt_textbox(
    slide: Any,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float,
    color: PptxRGBColor = PPT_INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _add_visual_panel(slide: Any, slide_data: dict[str, Any]) -> None:
    visual_items = [
        str(item) for item in slide_data.get("visual_items", []) if str(item).strip()
    ][:4]
    if not visual_items:
        visual_items = ["Conteúdo", "Aprendizagem", "Avaliação"]
    visual_title = str(slide_data.get("visual_title", "Relação pedagógica"))
    alt_text = str(slide_data.get("alt_text", visual_title))

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.45),
        Inches(1.55),
        Inches(6.1),
        Inches(4.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PPT_PALE
    panel.line.color.rgb = PptxRGBColor(211, 222, 232)
    _set_shape_alt_text(panel, alt_text)
    _add_ppt_textbox(
        slide,
        visual_title,
        6.8,
        1.83,
        5.4,
        0.45,
        size=22,
        color=PPT_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    item_count = len(visual_items)
    gap = 0.18
    usable_width = 5.35
    card_width = (usable_width - gap * (item_count - 1)) / item_count
    start_x = 6.83
    kind = str(slide_data.get("visual_kind", "conceito"))
    palette = (PPT_BLUE, PPT_TEAL, PPT_GOLD, PPT_NAVY)
    if kind == "processo":
        for index in range(item_count - 1):
            x = start_x + card_width * (index + 1) + gap * index + 0.02
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x),
                Inches(3.55),
                Inches(x + gap - 0.04),
                Inches(3.55),
            )
            connector.line.color.rgb = PPT_MUTED
            connector.line.width = Pt(2)

    for index, item in enumerate(visual_items):
        x = start_x + index * (card_width + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.65),
            Inches(card_width),
            Inches(1.8),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = PPT_WHITE
        card.line.color.rgb = palette[index]
        card.line.width = Pt(2.25)
        card.text_frame.clear()
        card.text_frame.word_wrap = True
        card.text_frame.margin_left = Inches(0.12)
        card.text_frame.margin_right = Inches(0.12)
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = card.text_frame.paragraphs[0]
        paragraph.text = item
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16 if len(item) < 45 else 14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PPT_INK
        _set_shape_alt_text(card, f"{index + 1}. {item}. {alt_text}")

        number = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + card_width / 2 - 0.21),
            Inches(2.4),
            Inches(0.42),
            Inches(0.42),
        )
        number.fill.solid()
        number.fill.fore_color.rgb = palette[index]
        number.line.color.rgb = palette[index]
        number.text_frame.clear()
        number.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        number_paragraph = number.text_frame.paragraphs[0]
        number_paragraph.text = str(index + 1)
        number_paragraph.alignment = PP_ALIGN.CENTER
        number_paragraph.font.size = Pt(12)
        number_paragraph.font.bold = True
        number_paragraph.font.color.rgb = PPT_WHITE

    _add_ppt_textbox(
        slide,
        kind.upper(),
        8.45,
        4.85,
        2.2,
        0.38,
        size=11,
        color=PPT_TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _add_raster_image_panel(
    slide: Any,
    slide_data: dict[str, Any],
    asset: dict[str, Any],
) -> bool:
    """Insere uma imagem documental ou gerada preservando proporções e acessibilidade."""

    encoded = str(asset.get("data_base64", "")).strip()
    if not encoded:
        return False
    try:
        blob = base64.b64decode(encoded, validate=True)
        width_px, height_px = PptxImage.from_blob(blob).size
    except Exception:
        return False
    if width_px <= 0 or height_px <= 0:
        return False

    panel_x, panel_y, panel_w, panel_h = 6.45, 1.55, 6.1, 4.95
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(panel_x),
        Inches(panel_y),
        Inches(panel_w),
        Inches(panel_h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PPT_PALE
    panel.line.color.rgb = PptxRGBColor(211, 222, 232)

    visual_title = str(slide_data.get("visual_title", "Imagem documental"))
    alt_text = str(slide_data.get("alt_text", visual_title))
    _add_ppt_textbox(
        slide,
        visual_title,
        6.8,
        1.78,
        5.4,
        0.5,
        size=21,
        color=PPT_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    max_w, max_h = 5.25, 3.65
    ratio = width_px / height_px
    draw_w = max_w
    draw_h = draw_w / ratio
    if draw_h > max_h:
        draw_h = max_h
        draw_w = draw_h * ratio
    draw_x = panel_x + (panel_w - draw_w) / 2
    draw_y = 2.48 + (max_h - draw_h) / 2
    picture = slide.shapes.add_picture(
        io.BytesIO(blob),
        Inches(draw_x),
        Inches(draw_y),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )
    _set_shape_alt_text(picture, alt_text)

    if asset.get("origin_type") == "ai_generated":
        provider = str(asset.get("provider", "IA")).strip()
        model = str(asset.get("model", "")).strip()
        provenance = "Gerada por IA · " + provider
        if model:
            provenance += f" · {model}"
    else:
        location = str(asset.get("source_location", "")).strip()
        filename = str(asset.get("source_file", "")).strip() or str(
            asset.get("filename", "")
        ).strip()
        provenance = filename
        if location:
            provenance += f" · {location}"
    if provenance:
        _add_ppt_textbox(
            slide,
            provenance,
            6.82,
            6.18,
            5.35,
            0.22,
            size=9,
            color=PPT_MUTED,
            align=PP_ALIGN.CENTER,
        )
    return True


def _add_slide_source(slide: Any, source: str, *, light: bool = False) -> None:
    _add_ppt_textbox(
        slide,
        "Fonte visual: " + (source or f"Diagrama nativo gerado pelo {APP_NAME}."),
        0.75,
        6.93,
        11.8,
        0.25,
        size=9,
        color=PptxRGBColor(208, 220, 230) if light else PPT_MUTED,
    )


def export_presentation(state: dict[str, Any], output_path: Path | str | None = None) -> str:
    """Cria uma apresentação visual, editável e pedagogicamente alinhada."""

    slides = state.get("resources", {}).get("presentation_outline", [])
    if not slides:
        raise ValueError("A sessão não contém uma apresentação para exportar.")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    visual_assets = {
        str(asset.get("id", "")): asset
        for collection in ("source_images", "generated_images")
        for asset in state.get(collection, [])
        if isinstance(asset, dict)
        and str(asset.get("id", "")).strip()
        and asset.get("approved") is True
    }

    for index, slide_data in enumerate(slides):
        slide = presentation.slides.add_slide(blank_layout)
        bullets = [str(item) for item in slide_data.get("bullets", []) if str(item).strip()]
        visual_items = [
            str(item) for item in slide_data.get("visual_items", []) if str(item).strip()
        ][:4]
        source = str(slide_data.get("visual_source", ""))
        alt_text = str(slide_data.get("alt_text", ""))
        is_title = index == 0
        is_closing = index == len(slides) - 1

        if is_title or is_closing:
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = PPT_NAVY
            _add_ppt_textbox(
                slide,
                APP_NAME,
                0.75,
                0.35,
                2.2,
                0.4,
                size=17,
                color=PPT_TEAL,
                bold=True,
            )
            _add_ppt_textbox(
                slide,
                slide_data["title"],
                0.85,
                1.15,
                11.65,
                1.35,
                size=50 if is_title else 40,
                color=PPT_WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            if bullets:
                _add_ppt_textbox(
                    slide,
                    bullets[0],
                    1.65,
                    2.5,
                    10.0,
                    0.65,
                    size=20,
                    color=PptxRGBColor(216, 228, 238),
                    align=PP_ALIGN.CENTER,
                )
            items = visual_items or bullets[1:4] or [
                "Taxonomia",
                "Programa da UC",
                "Recursos alinhados",
            ]
            items = items[:4]
            item_width = 2.65
            total_width = len(items) * item_width + (len(items) - 1) * 0.35
            start_x = (13.333 - total_width) / 2
            for item_index, item in enumerate(items):
                x = start_x + item_index * (item_width + 0.35)
                card = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(x),
                    Inches(4.15),
                    Inches(item_width),
                    Inches(1.15),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = (PPT_BLUE, PPT_TEAL, PPT_GOLD, PPT_BLUE)[item_index]
                card.line.color.rgb = card.fill.fore_color.rgb
                card.text_frame.clear()
                card.text_frame.word_wrap = True
                card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                paragraph = card.text_frame.paragraphs[0]
                paragraph.text = item
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(17)
                paragraph.font.bold = True
                paragraph.font.color.rgb = PPT_WHITE
                _set_shape_alt_text(card, alt_text or str(item))
            _add_ppt_textbox(
                slide,
                str(slide_data.get("visual_title", "Percurso pedagógico")),
                4.25,
                5.58,
                4.8,
                0.35,
                size=13,
                color=PptxRGBColor(205, 219, 230),
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            _add_slide_source(slide, source, light=True)
            continue

        top_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            presentation.slide_width,
            Inches(0.16),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PPT_TEAL
        top_bar.line.fill.background()
        _add_ppt_textbox(
            slide,
            slide_data["title"],
            0.7,
            0.43,
            11.9,
            0.72,
            size=35,
            color=PPT_NAVY,
            bold=True,
        )
        outcome_id = str(slide_data.get("outcome_id", "")).strip()
        if outcome_id:
            pill = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.78),
                Inches(1.4),
                Inches(1.0),
                Inches(0.42),
            )
            pill.fill.solid()
            pill.fill.fore_color.rgb = PPT_TEAL
            pill.line.color.rgb = PPT_TEAL
            pill.text_frame.clear()
            pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = pill.text_frame.paragraphs[0]
            paragraph.text = outcome_id
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPT_WHITE

        bullet_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.95), Inches(5.25), Inches(4.45)
        )
        bullet_frame = bullet_box.text_frame
        bullet_frame.clear()
        bullet_frame.word_wrap = True
        bullet_frame.margin_left = Inches(0.18)
        bullet_frame.margin_right = Inches(0.12)
        for bullet_index, bullet in enumerate(bullets):
            paragraph = (
                bullet_frame.paragraphs[0]
                if bullet_index == 0
                else bullet_frame.add_paragraph()
            )
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(18 if len(bullets) <= 4 else 16)
            paragraph.font.color.rgb = PPT_INK
            paragraph.space_after = Pt(13)
        visual_mode = str(slide_data.get("visual_mode", "diagrama"))
        visual_asset_id = str(slide_data.get("visual_asset_id", "")).strip()
        used_raster_image = False
        if visual_mode in {"documento", "ia"} and visual_asset_id:
            asset = visual_assets.get(visual_asset_id)
            if asset is not None:
                used_raster_image = _add_raster_image_panel(
                    slide, slide_data, asset
                )
        if not used_raster_image:
            _add_visual_panel(slide, slide_data)
            if visual_mode == "documento":
                source = (
                    f"Diagrama nativo gerado pelo {APP_NAME}; a imagem documental "
                    "selecionada não estava disponível no momento da exportação."
                )
            elif visual_mode == "ia":
                source = (
                    f"Diagrama nativo gerado pelo {APP_NAME}; a imagem gerada por IA "
                    "não estava aprovada ou disponível no momento da exportação."
                )
        _add_slide_source(slide, source)

    if output_path is None:
        with NamedTemporaryFile(prefix="coeria_", suffix=".pptx", delete=False) as temp_file:
            destination = Path(temp_file.name)
    else:
        destination = Path(output_path)
    presentation.save(destination)
    return str(destination)


def _export_worksheet(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["lesson_worksheet"]
    document = Document()
    _set_document_defaults(document)
    _add_document_title(document, data["title"], state["course"]["unit_name"])
    document.add_paragraph(data["overview"])
    document.add_heading("Instruções", level=1)
    document.add_paragraph(data["instructions"])
    for section in data["sections"]:
        document.add_heading(section["heading"], level=1)
        document.add_paragraph(section["content"])
        outcomes = ", ".join(section["outcome_ids"])
        if outcomes:
            document.add_paragraph(f"Resultados associados: {outcomes}")
        if section["activity"]:
            document.add_heading("Atividade", level=2)
            document.add_paragraph(section["activity"])
    document.save(output_path)


def _export_test(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["test"]
    document = Document()
    _set_document_defaults(document)
    _add_document_title(document, data["title"], state["course"]["unit_name"])
    document.add_paragraph(data["instructions"])
    document.add_paragraph(f"Cotação total: {data['total_points']} pontos")
    document.add_heading("Questões", level=1)
    for question in data["questions"]:
        document.add_heading(
            f"{question['id']} — {question['points']} pontos", level=2
        )
        document.add_paragraph(question["prompt"])
        document.add_paragraph(
            f"Tipo: {question['question_type']} · Resultado: {question['outcome_id']}"
        )
        document.add_paragraph("\n\n")
    document.add_page_break()
    document.add_heading("Chave de correção", level=1)
    for question in data["questions"]:
        document.add_heading(question["id"], level=2)
        document.add_paragraph(question["answer_key"])
    document.save(output_path)


def _export_practical_activity(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["practical_activity"]
    document = Document()
    _set_document_defaults(document)
    _add_document_title(document, data["title"], state["course"]["unit_name"])
    document.add_paragraph(data["context"])
    document.add_paragraph(f"Duração prevista: {data['duration_minutes']} minutos")
    document.add_heading("Materiais", level=1)
    for material in data["materials"]:
        document.add_paragraph(material, style="List Bullet")
    document.add_heading("Etapas", level=1)
    for step in sorted(data["steps"], key=lambda item: item["order"]):
        outcomes = ", ".join(step["outcome_ids"])
        document.add_paragraph(
            f"{step['order']}. {step['instruction']} ({outcomes})", style="List Number"
        )
    document.add_heading("Entregáveis", level=1)
    for deliverable in data["deliverables"]:
        document.add_paragraph(deliverable, style="List Bullet")
    document.add_heading("Critérios", level=1)
    for criterion in data["criteria"]:
        document.add_paragraph(
            f"{criterion['criterion']} ({criterion['weight']}%): {criterion['description']}"
        )
    document.save(output_path)


def _export_worksheet_latex(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["lesson_worksheet"]
    body = [
        _latex_escape(data.get("overview", "")),
        r"\section{Instruções}",
        _latex_escape(data.get("instructions", "")),
    ]
    for section in data.get("sections", []):
        body.extend(
            [
                rf"\section{{{_latex_escape(section.get('heading', ''))}}}",
                _latex_escape(section.get("content", "")),
            ]
        )
        outcomes = ", ".join(section.get("outcome_ids", []))
        if outcomes:
            body.append(rf"\textbf{{Resultados associados:}} {_latex_escape(outcomes)}")
        if section.get("activity"):
            body.extend(
                [
                    r"\subsection{Atividade}",
                    _latex_escape(section.get("activity", "")),
                ]
            )
    content = _latex_document(
        str(data.get("title", "Ficha de aula")),
        str(state.get("course", {}).get("unit_name", "")),
        body,
    )
    _write_latex_document(content, output_path, prefix="coeria_ficha_")


def _export_test_latex(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["test"]
    body = [
        _latex_escape(data.get("instructions", "")),
        rf"\textbf{{Cotação total:}} {_latex_escape(data.get('total_points', 0))} pontos",
        r"\section{Questões}",
    ]
    questions = data.get("questions", [])
    for question in questions:
        body.extend(
            [
                rf"\subsection{{{_latex_escape(question.get('id', ''))} --- {_latex_escape(question.get('points', 0))} pontos}}",
                _latex_escape(question.get("prompt", "")) + r"\par",
                rf"\textbf{{Tipo:}} {_latex_escape(question.get('question_type', ''))}\par",
                rf"\textbf{{Resultado associado:}} {_latex_escape(question.get('outcome_id', ''))}",
                r"\vspace{3\baselineskip}",
            ]
        )
    body.extend([r"\newpage", r"\section{Chave de correção}"])
    for question in questions:
        body.extend(
            [
                rf"\subsection{{{_latex_escape(question.get('id', ''))}}}",
                _latex_escape(question.get("answer_key", "")),
            ]
        )
    content = _latex_document(
        str(data.get("title", "Teste")),
        str(state.get("course", {}).get("unit_name", "")),
        body,
    )
    _write_latex_document(content, output_path, prefix="coeria_teste_")


def _export_practical_activity_latex(state: dict[str, Any], output_path: Path) -> None:
    data = state["resources"]["practical_activity"]
    body = [
        _latex_escape(data.get("context", "")),
        rf"\textbf{{Duração prevista:}} {_latex_escape(data.get('duration_minutes', 0))} minutos",
        r"\section{Materiais}",
        _latex_itemize(list(data.get("materials", []))),
        r"\section{Etapas}",
        _latex_itemize(
            [
                (
                    f"{step.get('instruction', '')} "
                    f"({', '.join(step.get('outcome_ids', []))})"
                )
                for step in sorted(
                    data.get("steps", []), key=lambda item: item.get("order", 0)
                )
            ],
            ordered=True,
        ),
        r"\section{Entregáveis}",
        _latex_itemize(list(data.get("deliverables", []))),
        r"\section{Critérios}",
        _latex_itemize(
            [
                (
                    f"{criterion.get('criterion', '')} "
                    f"({criterion.get('weight', 0)}%): "
                    f"{criterion.get('description', '')}"
                )
                for criterion in data.get("criteria", [])
            ]
        ),
    ]
    content = _latex_document(
        str(data.get("title", "Atividade prática")),
        str(state.get("course", {}).get("unit_name", "")),
        body,
    )
    _write_latex_document(content, output_path, prefix="coeria_atividade_")


def _csv_bytes(headers: list[str], rows: list[list[Any]]) -> bytes:
    stream = io.StringIO(newline="")
    writer = csv.writer(stream)
    writer.writerow(headers)
    writer.writerows(rows)
    return stream.getvalue().encode("utf-8-sig")


def _register_latex_document(
    generated: list[tuple[Path, str]],
    source: Path,
    compiled_pdfs: list[str],
) -> None:
    generated.append((source, source.name))
    pdf_path_text = compile_latex_pdf(source)
    if pdf_path_text:
        pdf_path = Path(pdf_path_text)
        generated.append((pdf_path, pdf_path.name))
        compiled_pdfs.append(pdf_path.name)


def export_resource_package(
    state: dict[str, Any],
    document_formats: list[str] | tuple[str, ...] | str | None = None,
) -> str:
    """Cria um ZIP com os recursos selecionados e os elementos de auditoria."""

    if state.get("status") != "completed":
        raise ValueError("A sessão deve estar concluída antes da exportação.")
    resources = attach_quality_report(state, state.get("resources", {}))
    quality = resources.get("quality", {})
    if not quality.get("passed"):
        raise ValueError("A validação automática detetou erros bloqueantes nos recursos.")

    formats = normalize_document_formats(document_formats)
    pdf_compilation_enabled = (
        DOCUMENT_FORMAT_LATEX in formats and latex_pdf_compilation_enabled()
    )
    latex_compiler_name: str | None = None
    if pdf_compilation_enabled:
        latex_compiler_name = Path(_latex_compiler_path()).name
    selected = set(resources.get("selected_types", []))
    course_stem = _safe_stem(state["course"]["unit_name"])

    with TemporaryDirectory(prefix="coeria_export_") as temporary_directory:
        temporary_path = Path(temporary_directory)
        generated: list[tuple[Path, str]] = []
        primary_products: list[str] = []
        compiled_pdfs: list[str] = []

        if DOCUMENT_FORMAT_WORD in formats:
            program_path = temporary_path / f"{course_stem}_programa_uc.docx"
            export_program_document(state, program_path)
            generated.append((program_path, program_path.name))
            primary_products.append(program_path.name)
        if DOCUMENT_FORMAT_LATEX in formats:
            program_latex_path = temporary_path / f"{course_stem}_programa_uc.tex"
            export_program_latex(state, program_latex_path)
            _register_latex_document(
                generated,
                program_latex_path,
                compiled_pdfs,
            )
            primary_products.append(program_latex_path.name)

        if RESOURCE_PRESENTATION in selected:
            path = temporary_path / f"{course_stem}_apresentacao.pptx"
            export_presentation(state, path)
            generated.append((path, path.name))
        if RESOURCE_WORKSHEET in selected:
            if DOCUMENT_FORMAT_WORD in formats:
                path = temporary_path / f"{course_stem}_ficha_aula.docx"
                _export_worksheet(state, path)
                generated.append((path, path.name))
            if DOCUMENT_FORMAT_LATEX in formats:
                path = temporary_path / f"{course_stem}_ficha_aula.tex"
                _export_worksheet_latex(state, path)
                _register_latex_document(generated, path, compiled_pdfs)
        if RESOURCE_TEST in selected:
            if DOCUMENT_FORMAT_WORD in formats:
                path = temporary_path / f"{course_stem}_teste.docx"
                _export_test(state, path)
                generated.append((path, path.name))
            if DOCUMENT_FORMAT_LATEX in formats:
                path = temporary_path / f"{course_stem}_teste.tex"
                _export_test_latex(state, path)
                _register_latex_document(generated, path, compiled_pdfs)
        if RESOURCE_PRACTICAL in selected:
            if DOCUMENT_FORMAT_WORD in formats:
                path = temporary_path / f"{course_stem}_atividade_pratica.docx"
                _export_practical_activity(state, path)
                generated.append((path, path.name))
            if DOCUMENT_FORMAT_LATEX in formats:
                path = temporary_path / f"{course_stem}_atividade_pratica.tex"
                _export_practical_activity_latex(state, path)
                _register_latex_document(generated, path, compiled_pdfs)

        alignment_csv_rows = [
            [
                row.get("outcome_id", ""),
                row.get("result", ""),
                ", ".join(row.get("content_ids", [])),
                row.get("taxonomy", ""),
                row.get("taxonomy_level", ""),
                ", ".join(row.get("assessment_ids", [])),
                ", ".join(row.get("assessment_purposes", [])),
                ", ".join(row.get("teaching_activity_ids", [])),
                row.get("status", ""),
                row.get("rationale", ""),
            ]
            for row in derive_alignment_rows(state)
        ]
        audit_rows = [
            [item.get("timestamp", ""), item.get("stage", ""), item.get("event", ""), item.get("feedback", "")]
            for item in state.get("audit", [])
        ]
        manifest = {
            "application": APP_NAME,
            "application_name": APP_FULL_NAME,
            "session_id": state.get("session_id", ""),
            "ai_provider": state.get("ai_provider", "OpenAI"),
            "course": state.get("course", {}),
            "taxonomy": state.get("course", {}).get("taxonomy_type", "SOLO"),
            "selected_resources": list(resources.get("selected_types", [])),
            "document_formats": list(formats),
            "latex_pdf_compilation": {
                "enabled": pdf_compilation_enabled,
                "compiler": latex_compiler_name,
                "generated_files": compiled_pdfs,
            },
            "quality": quality,
            "primary_product": primary_products[0],
            "primary_products": primary_products,
            "files": [archive_name for _path, archive_name in generated],
            "visual_assets": {
                "document_images": [
                    {
                        **{
                            key: asset.get(key)
                            for key in (
                                "id", "origin_type", "candidate_kind", "source_file",
                                "source_location", "filename", "media_type", "width_px",
                                "height_px", "image_mode", "approved",
                            )
                        },
                        "available_to_llm": asset.get("origin_type")
                        != "user_uploaded",
                        "used_in_presentation": str(asset.get("id", ""))
                        in {
                            str(slide.get("visual_asset_id", ""))
                            for slide in resources.get("presentation_outline", [])
                            if isinstance(slide, dict)
                        },
                    }
                    for asset in state.get("source_images", [])
                    if isinstance(asset, dict)
                ],
                "ai_generated_images": [
                    {
                        key: asset.get(key)
                        for key in (
                            "id", "origin_type", "provider", "model", "prompt",
                            "size", "quality", "output_format", "filename", "media_type",
                            "alt_text", "approved", "created_at",
                        )
                    }
                    for asset in state.get("generated_images", [])
                    if isinstance(asset, dict)
                ],
            },
        }

        with NamedTemporaryFile(
            prefix=f"coeria_{course_stem}_", suffix=".zip", delete=False
        ) as temp_file:
            package_path = Path(temp_file.name)
        with zipfile.ZipFile(package_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for path, archive_name in generated:
                archive.write(path, archive_name)
            archive.writestr(
                "sintese_alinhamento.csv",
                _csv_bytes(
                    [
                        "Resultado", "Descrição", "Conteúdos",
                        "Taxonomia", "Nível", "Avaliações", "Finalidade",
                        "Atividades de ensino-aprendizagem", "Estado",
                        "Justificação"
                    ],
                    alignment_csv_rows,
                ),
            )
            archive.writestr(
                "rastreabilidade.csv",
                _csv_bytes(["Data", "Etapa", "Evento", "Feedback"], audit_rows),
            )
            archive.writestr(
                "manifesto.json",
                json.dumps(manifest, ensure_ascii=False, indent=2).encode("utf-8"),
            )
            archive.writestr(
                "estado_sessao.json",
                json.dumps(state, ensure_ascii=False, indent=2).encode("utf-8"),
            )

    return str(package_path)
