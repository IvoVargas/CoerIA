"""Extração segura de texto dos ficheiros fornecidos pelo docente."""

from __future__ import annotations

import base64
import re
from hashlib import sha256
from pathlib import Path
from typing import Any, Iterable
from zipfile import BadZipFile, ZipFile

from .branding import config_value
from .image_utils import ImageValidationError, build_thumbnail, normalise_image_bytes


SUPPORTED_SOURCE_SUFFIXES = {".txt", ".md", ".tex", ".pdf", ".docx", ".pptx"}
DIRECT_SOURCE_MARKER = "[Texto introduzido pelo docente]"
DEFAULT_MAX_SOURCE_CHARS = 120_000
DEFAULT_MAX_RAW_SOURCE_CHARS = 2_000_000
DEFAULT_MAX_FILE_BYTES = 50 * 1024 * 1024
DEFAULT_MAX_TOTAL_UPLOAD_BYTES = 100 * 1024 * 1024
DEFAULT_MAX_EXTRACTED_IMAGE_BYTES = 5 * 1024 * 1024
DEFAULT_MAX_EXTRACTED_IMAGES = 30
DEFAULT_MIN_EXTRACTED_IMAGE_WIDTH = 240
DEFAULT_MIN_EXTRACTED_IMAGE_HEIGHT = 160
DEFAULT_MIN_IMAGE_FOOTPRINT_WIDTH_PT = 72
DEFAULT_MIN_IMAGE_FOOTPRINT_HEIGHT_PT = 54
DEFAULT_PDF_COMPOSITE_GAP_PT = 18
DEFAULT_PDF_RENDER_SCALE = 2.0


class SourceIngestionError(ValueError):
    """Erro compreensível durante a leitura de uma fonte documental."""


def _configured_limit(suffix: str, default: int) -> int:
    name = f"COERIA_{suffix}"
    raw = config_value(suffix, str(default))
    try:
        value = int(raw)
    except ValueError as error:
        raise SourceIngestionError(f"A variável {name} deve ser um número inteiro.") from error
    if value <= 0:
        raise SourceIngestionError(f"A variável {name} deve ser superior a zero.")
    return value


def configured_max_file_bytes() -> int:
    """Limite por ficheiro partilhado pela interface e pela ingestão."""

    return _configured_limit("MAX_FILE_BYTES", DEFAULT_MAX_FILE_BYTES)


def configured_max_total_upload_bytes() -> int:
    """Orçamento acumulado dos ficheiros escolhidos antes de iniciar a sessão."""

    return _configured_limit(
        "MAX_TOTAL_UPLOAD_BYTES", DEFAULT_MAX_TOTAL_UPLOAD_BYTES
    )


def _read_plain_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise SourceIngestionError(
            "A leitura de PDF requer a dependência pypdf. Reinstale os requisitos."
        ) from error

    try:
        reader = PdfReader(str(path))
        return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages)
    except Exception as error:
        raise SourceIngestionError(f"Não foi possível extrair texto de {path.name}.") from error


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as error:
        raise SourceIngestionError(
            "A leitura de DOCX requer a dependência python-docx. Reinstale os requisitos."
        ) from error

    try:
        document = Document(str(path))
        parts = [paragraph.text.strip() for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(part for part in parts if part)
    except Exception as error:
        raise SourceIngestionError(f"Não foi possível extrair texto de {path.name}.") from error


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise SourceIngestionError(
            "A leitura de PPTX requer a dependência python-pptx. Reinstale os requisitos."
        ) from error

    try:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
            if any(slide_text):
                parts.append(f"Slide {index}\n" + "\n".join(text for text in slide_text if text))
        return "\n\n".join(parts)
    except Exception as error:
        raise SourceIngestionError(f"Não foi possível extrair texto de {path.name}.") from error


def _candidate_limit() -> int:
    return _configured_limit("MAX_EXTRACTED_IMAGES", DEFAULT_MAX_EXTRACTED_IMAGES)


def _source_image_asset(
    *,
    data: bytes,
    source_file: str,
    filename: str,
    location: str = "",
    candidate_kind: str = "embedded",
    page_number: int | None = None,
    source_bbox: tuple[float, float, float, float] | None = None,
    object_count: int = 1,
) -> dict[str, Any] | None:
    """Valida e normaliza uma imagem documental antes de a catalogar."""

    max_bytes = _configured_limit(
        "MAX_EXTRACTED_IMAGE_BYTES", DEFAULT_MAX_EXTRACTED_IMAGE_BYTES
    )
    try:
        normalized = normalise_image_bytes(
            data,
            filename=filename,
            max_bytes=max_bytes,
        )
    except ImageValidationError:
        return None

    min_width = _configured_limit(
        "MIN_EXTRACTED_IMAGE_WIDTH", DEFAULT_MIN_EXTRACTED_IMAGE_WIDTH
    )
    min_height = _configured_limit(
        "MIN_EXTRACTED_IMAGE_HEIGHT", DEFAULT_MIN_EXTRACTED_IMAGE_HEIGHT
    )
    if (
        int(normalized["width_px"]) < min_width
        or int(normalized["height_px"]) < min_height
    ):
        return None

    normalized_data = bytes(normalized["data"])
    digest = sha256(
        source_file.encode("utf-8")
        + b"\0"
        + location.encode("utf-8")
        + b"\0"
        + candidate_kind.encode("utf-8")
        + b"\0"
        + normalized_data
    ).hexdigest()[:20]
    thumbnail = build_thumbnail(normalized_data)

    asset: dict[str, Any] = {
        "id": f"document-{digest}",
        "origin_type": "document",
        "candidate_kind": candidate_kind,
        "source_file": source_file,
        "source_location": location,
        "filename": str(normalized["filename"]),
        "media_type": str(normalized["media_type"]),
        "data_base64": base64.b64encode(normalized_data).decode("ascii"),
        "width_px": int(normalized["width_px"]),
        "height_px": int(normalized["height_px"]),
        "image_mode": "RGB",
        "alt_text": "",
        "approved": False,
        **thumbnail,
    }
    if page_number is not None:
        asset["page_number"] = int(page_number)
    if source_bbox is not None:
        asset["source_bbox"] = [round(float(value), 2) for value in source_bbox]
    if object_count > 1:
        asset["composite_object_count"] = int(object_count)
    return asset


def _bbox_union(
    boxes: list[tuple[float, float, float, float]],
) -> tuple[float, float, float, float]:
    return (
        min(box[0] for box in boxes),
        min(box[1] for box in boxes),
        max(box[2] for box in boxes),
        max(box[3] for box in boxes),
    )


def _boxes_near(
    first: tuple[float, float, float, float],
    second: tuple[float, float, float, float],
    gap: float,
) -> bool:
    return not (
        first[2] + gap < second[0]
        or second[2] + gap < first[0]
        or first[3] + gap < second[1]
        or second[3] + gap < first[1]
    )


def _cluster_pdf_fragments(
    entries: list[dict[str, Any]],
    gap: float,
) -> list[list[dict[str, Any]]]:
    """Agrupa objetos raster próximos antes de extrair fragmentos isolados."""

    remaining = list(entries)
    clusters: list[list[dict[str, Any]]] = []
    while remaining:
        cluster = [remaining.pop(0)]
        changed = True
        while changed:
            changed = False
            union = _bbox_union([item["bbox"] for item in cluster])
            for item in list(remaining):
                if _boxes_near(union, item["bbox"], gap):
                    cluster.append(item)
                    remaining.remove(item)
                    changed = True
        clusters.append(cluster)
    return clusters


def _balanced_candidates(
    grouped: dict[Any, list[dict[str, Any]]],
    limit: int,
) -> list[dict[str, Any]]:
    """Seleciona em rondas para impedir uma página/ficheiro de monopolizar o catálogo."""

    if limit <= 0:
        return []
    ordered_keys = sorted(grouped, key=lambda value: str(value))
    queues = {
        key: sorted(
            grouped[key],
            key=lambda asset: (
                0 if asset.get("candidate_kind") == "composite_render" else 1,
                -(int(asset.get("width_px", 0)) * int(asset.get("height_px", 0))),
                str(asset.get("id", "")),
            ),
        )
        for key in ordered_keys
    }
    selected: list[dict[str, Any]] = []
    while len(selected) < limit and any(queues[key] for key in ordered_keys):
        for key in ordered_keys:
            if queues[key] and len(selected) < limit:
                selected.append(queues[key].pop(0))
    return selected


def _extract_pdf_images(path: Path) -> list[dict[str, Any]]:
    """Examina o PDF completo e só depois aplica o limite ao catálogo visual."""

    try:
        import pymupdf
    except ImportError as error:
        raise SourceIngestionError(
            "A extração visual de PDF requer a dependência PyMuPDF. Reinstale os requisitos."
        ) from error

    per_page: dict[int, list[dict[str, Any]]] = {}
    try:
        document = pymupdf.open(str(path))
    except Exception:
        return []

    try:
        gap = float(config_value("PDF_COMPOSITE_GAP_PT", str(DEFAULT_PDF_COMPOSITE_GAP_PT)))
        render_scale = float(config_value("PDF_RENDER_SCALE", str(DEFAULT_PDF_RENDER_SCALE)))
        min_footprint_w = float(
            config_value(
                "MIN_IMAGE_FOOTPRINT_WIDTH_PT",
                str(DEFAULT_MIN_IMAGE_FOOTPRINT_WIDTH_PT),
            )
        )
        min_footprint_h = float(
            config_value(
                "MIN_IMAGE_FOOTPRINT_HEIGHT_PT",
                str(DEFAULT_MIN_IMAGE_FOOTPRINT_HEIGHT_PT),
            )
        )

        # Deliberadamente percorre todas as páginas antes de qualquer corte global.
        for page_index in range(document.page_count):
            page = document.load_page(page_index)
            page_number = page_index + 1
            page_area = max(float(page.rect.width * page.rect.height), 1.0)
            fragments: list[dict[str, Any]] = []
            extracted_by_xref: dict[int, tuple[bytes, str]] = {}

            for image_info in page.get_images(full=True):
                xref = int(image_info[0])
                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    rects = []
                if not rects:
                    continue
                if xref not in extracted_by_xref:
                    try:
                        extracted = document.extract_image(xref)
                        extracted_by_xref[xref] = (
                            bytes(extracted.get("image", b"")),
                            f"imagem-{xref}.{extracted.get('ext', 'bin')}",
                        )
                    except Exception:
                        extracted_by_xref[xref] = (b"", f"imagem-{xref}.bin")
                for instance_index, rect in enumerate(rects, start=1):
                    bbox = (
                        float(rect.x0),
                        float(rect.y0),
                        float(rect.x1),
                        float(rect.y1),
                    )
                    fragments.append(
                        {
                            "key": (xref, instance_index, bbox),
                            "xref": xref,
                            "bbox": bbox,
                            "data": extracted_by_xref[xref][0],
                            "filename": extracted_by_xref[xref][1],
                        }
                    )

            page_candidates: list[dict[str, Any]] = []
            suppressed: set[tuple[Any, ...]] = set()
            for cluster in _cluster_pdf_fragments(fragments, gap):
                if len(cluster) < 2:
                    continue
                union = _bbox_union([item["bbox"] for item in cluster])
                width_pt = union[2] - union[0]
                height_pt = union[3] - union[1]
                coverage = max(width_pt * height_pt, 0.0) / page_area
                if (
                    width_pt < min_footprint_w
                    or height_pt < min_footprint_h
                    or coverage < 0.02
                    or coverage > 0.75
                ):
                    continue

                margin = max(8.0, min(18.0, gap))
                clip = pymupdf.Rect(
                    max(page.rect.x0, union[0] - margin),
                    max(page.rect.y0, union[1] - margin),
                    min(page.rect.x1, union[2] + margin),
                    min(page.rect.y1, union[3] + margin),
                )
                try:
                    pixmap = page.get_pixmap(
                        matrix=pymupdf.Matrix(render_scale, render_scale),
                        clip=clip,
                        alpha=False,
                    )
                    rendered = pixmap.tobytes("png")
                except Exception:
                    continue
                asset = _source_image_asset(
                    data=rendered,
                    source_file=path.name,
                    filename=f"pagina-{page_number}-figura-composta.png",
                    location=f"Página {page_number}",
                    candidate_kind="composite_render",
                    page_number=page_number,
                    source_bbox=(clip.x0, clip.y0, clip.x1, clip.y1),
                    object_count=len(cluster),
                )
                if asset is not None:
                    page_candidates.append(asset)
                    suppressed.update(item["key"] for item in cluster)

            for fragment in fragments:
                if fragment["key"] in suppressed:
                    continue
                x0, y0, x1, y1 = fragment["bbox"]
                if (x1 - x0) < min_footprint_w or (y1 - y0) < min_footprint_h:
                    continue
                asset = _source_image_asset(
                    data=fragment["data"],
                    source_file=path.name,
                    filename=fragment["filename"],
                    location=f"Página {page_number}",
                    candidate_kind="embedded",
                    page_number=page_number,
                    source_bbox=fragment["bbox"],
                )
                if asset is not None:
                    page_candidates.append(asset)

            if page_candidates:
                per_page[page_number] = page_candidates

        return _balanced_candidates(per_page, _candidate_limit())
    except Exception:
        return []
    finally:
        document.close()


def _extract_docx_images(path: Path) -> list[dict[str, Any]]:
    candidates: list[dict[str, Any]] = []
    try:
        with ZipFile(path) as archive:
            media_names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("word/media/") and not name.endswith("/")
            )
            # Examina todas as imagens do pacote antes de limitar o catálogo.
            for media_name in media_names:
                asset = _source_image_asset(
                    data=archive.read(media_name),
                    source_file=path.name,
                    filename=Path(media_name).name,
                    candidate_kind="embedded",
                )
                if asset is not None:
                    candidates.append(asset)
        return _balanced_candidates({"documento": candidates}, _candidate_limit())
    except (BadZipFile, KeyError, OSError):
        return []


def _extract_pptx_images(path: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as error:
        raise SourceIngestionError(
            "A extração de imagens de PPTX requer a dependência python-pptx. "
            "Reinstale os requisitos."
        ) from error

    per_slide: dict[int, list[dict[str, Any]]] = {}
    try:
        presentation = Presentation(str(path))

        def visit_shapes(shapes: Any, slide_index: int, bucket: list[dict[str, Any]]) -> None:
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    visit_shapes(shape.shapes, slide_index, bucket)
                    continue
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    image = shape.image
                    asset = _source_image_asset(
                        data=bytes(image.blob),
                        source_file=path.name,
                        filename=str(image.filename or f"imagem.{image.ext}"),
                        location=f"Slide {slide_index}",
                        candidate_kind="embedded",
                        page_number=slide_index,
                    )
                except Exception:
                    asset = None
                if asset is not None:
                    bucket.append(asset)

        # Percorre todos os slides antes de aplicar o limite.
        for slide_index, slide in enumerate(presentation.slides, start=1):
            bucket: list[dict[str, Any]] = []
            visit_shapes(slide.shapes, slide_index, bucket)
            if bucket:
                per_slide[slide_index] = bucket
        return _balanced_candidates(per_slide, _candidate_limit())
    except Exception:
        return []


def extract_file_images(
    path_value: str | Path,
) -> list[dict[str, Any]]:
    """Extrai candidatos visuais normalizados com proveniência rastreável."""

    path = Path(path_value)
    if not path.is_file():
        raise SourceIngestionError(
            f"O ficheiro {path.name or path} não está disponível."
        )

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_images(path)
    if suffix == ".docx":
        return _extract_docx_images(path)
    if suffix == ".pptx":
        return _extract_pptx_images(path)
    return []


def extract_source_images(
    file_paths: str | Path | Iterable[str | Path] | None = None,
) -> list[dict[str, Any]]:
    """Examina todas as fontes e equilibra o catálogo final entre ficheiros."""

    grouped: dict[str, list[dict[str, Any]]] = {}
    seen_ids: set[str] = set()
    for path_value in _normalise_paths(file_paths):
        path = Path(path_value)
        file_assets: list[dict[str, Any]] = []
        for asset in extract_file_images(path):
            asset_id = str(asset.get("id", ""))
            if not asset_id or asset_id in seen_ids:
                continue
            seen_ids.add(asset_id)
            file_assets.append(asset)
        if file_assets:
            grouped[path.name] = file_assets
    return _balanced_candidates(grouped, _candidate_limit())

def extract_file_text(path_value: str | Path) -> str:
    path = Path(path_value)
    if not path.is_file():
        raise SourceIngestionError(f"O ficheiro {path.name or path} não está disponível.")
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SOURCE_SUFFIXES:
        formats = ", ".join(sorted(SUPPORTED_SOURCE_SUFFIXES))
        raise SourceIngestionError(f"Formato não suportado. Utilize: {formats}.")

    max_bytes = configured_max_file_bytes()
    if path.stat().st_size > max_bytes:
        raise SourceIngestionError(
            f"O ficheiro {path.name} excede o limite de {max_bytes // (1024 * 1024)} MB."
        )

    if suffix in {".txt", ".md", ".tex"}:
        text = _read_plain_text(path)
    elif suffix == ".pdf":
        text = _read_pdf(path)
    elif suffix == ".docx":
        text = _read_docx(path)
    else:
        text = _read_pptx(path)

    clean = text.strip()
    if not clean:
        raise SourceIngestionError(
            f"Não foi encontrado texto utilizável em {path.name}. PDFs digitalizados podem exigir OCR."
        )
    return clean


def _normalise_paths(file_paths: str | Path | Iterable[str | Path] | None) -> list[str | Path]:
    if file_paths is None:
        return []
    if isinstance(file_paths, (str, Path)):
        return [file_paths]
    return list(file_paths)


def _combine_source_text(
    source_text: str,
    file_paths: str | Path | Iterable[str | Path] | None = None,
) -> str:
    parts: list[str] = []
    direct = (source_text or "").strip()
    if direct:
        parts.append(DIRECT_SOURCE_MARKER + "\n" + direct)

    for path_value in _normalise_paths(file_paths):
        path = Path(path_value)
        parts.append(f"[Ficheiro: {path.name}]\n{extract_file_text(path)}")

    return "\n\n".join(parts).strip()


def build_raw_source_text(
    source_text: str,
    file_paths: str | Path | Iterable[str | Path] | None = None,
) -> str:
    """Combina fontes extensas antes da eventual redução automática por IA."""

    combined = _combine_source_text(source_text, file_paths)
    max_chars = _configured_limit(
        "MAX_RAW_SOURCE_CHARS", DEFAULT_MAX_RAW_SOURCE_CHARS
    )
    if len(combined) > max_chars:
        raise SourceIngestionError(
            f"As fontes contêm {len(combined):,} caracteres e excedem o limite absoluto "
            f"de ingestão de {max_chars:,}. Remova documentos redundantes ou aumente "
            "COERIA_MAX_RAW_SOURCE_CHARS de forma consciente."
        )
    return combined


def build_source_text(
    source_text: str,
    file_paths: str | Path | Iterable[str | Path] | None = None,
) -> str:
    """Combina fontes respeitando o orçamento normal do contexto do pipeline.

    Mantém-se como API compatível para validações/testes. A aplicação usa
    :func:`build_raw_source_text` e reduz automaticamente fontes extensas antes
    de iniciar o fluxo pedagógico.
    """

    combined = build_raw_source_text(source_text, file_paths)
    max_chars = _configured_limit("MAX_SOURCE_CHARS", DEFAULT_MAX_SOURCE_CHARS)
    if len(combined) > max_chars:
        raise SourceIngestionError(
            f"As fontes contêm {len(combined):,} caracteres e excedem o limite de "
            f"{max_chars:,}."
        )
    return combined


def recover_direct_source_text(combined_text: str) -> str:
    """Recupera o texto direto de estados antigos com etiquetas de proveniência.

    Sessões novas guardam o valor original separadamente. Esta função existe
    para que sessões criadas antes dessa alteração não apresentem as etiquetas
    internas nem o conteúdo extraído dos ficheiros no campo do docente.
    """

    text = (combined_text or "").strip()
    if not text.startswith(DIRECT_SOURCE_MARKER):
        return "" if text.startswith("[Ficheiro: ") else text

    direct = text[len(DIRECT_SOURCE_MARKER) :].lstrip("\r\n")
    file_marker = re.search(r"\r?\n\r?\n\[Ficheiro: [^\]]+\]\r?\n", direct)
    if file_marker:
        direct = direct[: file_marker.start()]
    return direct.strip()
