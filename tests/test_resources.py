import base64
import json
import shutil
import unittest
import zipfile
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from os import environ
from tempfile import TemporaryDirectory
from unittest.mock import patch

from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from prism.agents import (
    AgentGenerationError,
    GenerationResult,
    _schema_for,
    validate_artifact,
)
from prism.application_service import ApplicationService
from prism.exporter import (
    _latex_itemize,
    compile_latex_pdf,
    export_presentation,
    export_program_document,
    export_program_latex,
    export_resource_package,
)
from prism.models import (
    CourseInput,
    RESOURCE_PRACTICAL,
    RESOURCE_PRESENTATION,
    RESOURCE_TEST,
    RESOURCE_WORKSHEET,
    SUPPORTED_RESOURCE_TYPES,
)
from prism.persistence import SQLiteSessionStore
from prism.presentation import render_current_artifact, render_resource_detail_sections
from prism.quality import PRESENTATION_ASSESSMENT_TITLE, evaluate_quality
from prism.workflow import (
    build_final_validation,
    create_session,
    create_test_agent,
    navigate_to_stage,
    review_current_stage,
)


class ResourceGenerationTests(unittest.TestCase):
    def setUp(self) -> None:
        self.course = CourseInput.create(
            unit_name="Programação Aplicada",
            source_text=(
                "Algoritmos e estruturas de controlo. Funções e modularidade. "
                "Testes automatizados e resolução de problemas aplicados."
            ),
            audience="Licenciatura",
            duration_hours=20,
            isced_f_code="0613",
            isced_f_name="Desenvolvimento e análise de software e aplicações",
            bibliography=(
                "Biggs, J., & Tang, C. (2011). Teaching for Quality Learning at University.\n"
                "Anderson, L. W., & Krathwohl, D. R. (2001). A Taxonomy for Learning."
            ),
        )
        self.agent = create_test_agent()

    def _resource_state(self):
        state = create_session(
            self.course,
            resource_types=list(SUPPORTED_RESOURCE_TYPES),
            agent=self.agent,
        )
        for _ in range(5):
            state = review_current_stage(state, "approve", agent=self.agent)
        self.assertEqual(state["current_stage"], "resources")
        return state

    def test_every_selected_resource_is_generated_and_validated(self) -> None:
        state = self._resource_state()
        resources = state["resources"]
        self.assertTrue(resources["presentation_outline"])
        self.assertTrue(
            all(
                slide.get("visual_mode") == "diagrama"
                and slide.get("visual_asset_id") == ""
                and slide.get("visual_prompt") == ""
                and slide.get("visual_kind")
                and slide.get("visual_title")
                and 2 <= len(slide.get("visual_items", [])) <= 4
                and slide.get("visual_source")
                and slide.get("alt_text")
                for slide in resources["presentation_outline"]
            )
        )
        assessment_slides = [
            slide
            for slide in resources["presentation_outline"]
            if slide["title"].startswith(PRESENTATION_ASSESSMENT_TITLE)
        ]
        self.assertTrue(assessment_slides)
        assessment_text = " ".join(
            str(text)
            for slide in assessment_slides
            for text in [slide["title"], *slide["bullets"]]
        )
        self.assertTrue(
            all(
                task["id"] in assessment_text
                and task["assessment_purpose"] in assessment_text
                and task["criterion"] in assessment_text
                for task in state["assessment_activities"]
            )
        )
        self.assertTrue(resources["lesson_worksheet"]["sections"])
        self.assertTrue(resources["test"]["questions"])
        self.assertTrue(resources["practical_activity"]["steps"])
        self.assertTrue(resources["quality"]["passed"])
        self.assertEqual(resources["quality"]["status"], "OK")

    def test_resource_detail_sections_include_only_selected_resources(self) -> None:
        state = self._resource_state()
        resources = deepcopy(state["resources"])
        resources["selected_types"] = [RESOURCE_TEST]

        sections = render_resource_detail_sections(resources)

        self.assertEqual([item["id"] for item in sections], ["test"])
        self.assertIn("Questões", sections[0]["content"])
        self.assertNotIn("Enunciado", sections[0]["content"])
        self.assertIn("Chave de correção", sections[0]["content"])
        self.assertNotIn("Atividade prática", sections[0]["content"])

    def test_presentation_detail_omits_the_redundant_source_column(self) -> None:
        state = self._resource_state()
        resources = deepcopy(state["resources"])
        resources["selected_types"] = [RESOURCE_PRESENTATION]

        sections = render_resource_detail_sections(resources)

        self.assertEqual([item["id"] for item in sections], ["presentation"])
        self.assertIn("| Modo visual |", sections[0]["content"])
        self.assertNotIn("| Fonte |", sections[0]["content"])

    def test_resource_quality_is_presented_only_in_final_validation(self) -> None:
        state = self._resource_state()

        resource_view = render_current_artifact(state)
        self.assertNotIn("Validação automática", resource_view)

        state = review_current_stage(state, "approve", agent=self.agent)
        self.assertEqual(state["current_stage"], "final_validation")
        final_view = render_current_artifact(state)
        self.assertIn("Qualidade automática dos recursos", final_view)
        self.assertIn("Correspondência dos recursos selecionados", final_view)
        self.assertIn("A seleção foi respeitada.", final_view)
        self.assertEqual(
            state["final_validation"]["resource_quality_checks"],
            state["resources"]["quality"]["checks"],
        )
        self.assertNotIn(
            "Os recursos selecionados devem cumprir os controlos determinísticos.",
            final_view,
        )

    def test_final_validation_refreshes_the_cached_resource_quality(self) -> None:
        state = self._resource_state()
        state["resources"]["quality"] = {
            "passed": False,
            "checks": [],
            "summary": {"errors": 1},
        }

        validated = navigate_to_stage(state, "final_validation")

        self.assertTrue(validated["resources"]["quality"]["passed"])
        self.assertTrue(validated["final_validation"]["passed"])

    def test_export_recomputes_a_stale_resource_quality_cache(self) -> None:
        state = review_current_stage(self._resource_state(), "approve", agent=self.agent)
        state = review_current_stage(state, "approve", agent=self.agent)
        state["resources"]["quality"] = {
            "passed": False,
            "checks": [],
            "summary": {"errors": 1},
        }

        package_path = Path(export_resource_package(state))
        try:
            self.assertTrue(package_path.is_file())
        finally:
            package_path.unlink(missing_ok=True)

    def test_document_image_is_approved_and_embedded_in_presentation(self) -> None:
        state = self._resource_state()
        image_buffer = BytesIO()
        Image.new("RGB", (320, 180), (30, 110, 170)).save(
            image_buffer, format="PNG"
        )
        asset_id = "document-test-image"
        state["source_images"] = [
            {
                "id": asset_id,
                "origin_type": "document",
                "source_file": "apoio.pdf",
                "source_location": "Página 3",
                "filename": "figura.png",
                "media_type": "image/png",
                "data_base64": base64.b64encode(image_buffer.getvalue()).decode(
                    "ascii"
                ),
                "alt_text": "",
                "approved": False,
            }
        ]
        slide = state["resources"]["presentation_outline"][1]
        slide["visual_mode"] = "documento"
        slide["visual_asset_id"] = asset_id
        slide["visual_source"] = "Imagem extraída de apoio.pdf, Página 3."
        slide["alt_text"] = "Figura documental usada para apoiar o conteúdo do slide."
        state["resources"]["quality"] = evaluate_quality(
            state, state["resources"]
        )
        self.assertTrue(state["resources"]["quality"]["passed"])

        state = review_current_stage(state, "approve", agent=self.agent)
        self.assertTrue(state["source_images"][0]["approved"])
        state = review_current_stage(state, "approve", agent=self.agent)
        self.assertEqual(state["status"], "completed")

        presentation_path = Path(export_presentation(state))
        try:
            presentation = Presentation(presentation_path)
            pictures = [
                shape
                for shape in presentation.slides[1].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            self.assertEqual(len(pictures), 1)
            descriptions = [
                properties.get("descr", "")
                for properties in pictures[0]._element.xpath(".//p:cNvPr")
            ]
            self.assertTrue(any("Figura documental" in item for item in descriptions))
        finally:
            presentation_path.unlink(missing_ok=True)

    def test_uploaded_image_is_valid_without_document_selection(self) -> None:
        state = self._resource_state()
        image_buffer = BytesIO()
        Image.new("RGB", (320, 180), (45, 125, 85)).save(
            image_buffer, format="PNG"
        )
        asset_id = "upload-test-image"
        state["source_images"] = [
            {
                "id": asset_id,
                "origin_type": "user_uploaded",
                "source_file": "imagem-local.png",
                "source_location": "Carregada pelo docente",
                "filename": "imagem-local.png",
                "media_type": "image/png",
                "data_base64": base64.b64encode(image_buffer.getvalue()).decode(
                    "ascii"
                ),
                "alt_text": "Imagem local usada para explicar o conteúdo.",
                "approved": False,
            }
        ]
        slide = state["resources"]["presentation_outline"][1]
        slide["visual_mode"] = "documento"
        slide["visual_asset_id"] = asset_id
        slide["visual_source"] = "Imagem fornecida pelo docente — imagem-local.png."
        slide["alt_text"] = "Imagem local usada para explicar o conteúdo."

        quality = evaluate_quality(state, state["resources"])

        self.assertTrue(quality["passed"])

    def test_manual_editor_image_generation_has_an_independent_limit(self) -> None:
        state = self._resource_state()
        slide = state["resources"]["presentation_outline"][1]

        class FakeGenerator:
            def generate(self, *, prompt, slide_number, alt_text):
                return {
                    "id": f"ai-manual-{slide_number}",
                    "origin_type": "ai_generated",
                    "provider": "Fornecedor de teste",
                    "model": "modelo-teste",
                    "prompt": prompt,
                    "alt_text": alt_text,
                    "approved": False,
                }

        with TemporaryDirectory() as temporary_directory, patch.dict(
            environ,
            {"COERIA_OPENAI_IMAGE_MAX_ADDITIONAL_EDITOR": "1"},
            clear=False,
        ):
            service = ApplicationService(
                SQLiteSessionStore(Path(temporary_directory) / "manual-images.db")
            )
            updated, asset = service.generate_presentation_editor_image(
                state,
                slide,
                2,
                "Representar visualmente o processo sem texto.",
                generator=FakeGenerator(),
            )
            self.assertEqual(asset["generation_mode"], "manual_editor")
            self.assertEqual(len(updated["generated_images"]), 1)
            with self.assertRaisesRegex(ValueError, "limite de 1 imagem"):
                service.generate_presentation_editor_image(
                    updated,
                    slide,
                    2,
                    "Gerar uma segunda imagem.",
                    generator=FakeGenerator(),
                )

    def test_ai_generated_image_is_approved_and_embedded_in_presentation(self) -> None:
        state = self._resource_state()
        image_buffer = BytesIO()
        Image.new("RGB", (320, 180), (90, 50, 140)).save(
            image_buffer, format="PNG"
        )
        asset_id = "ai-test-image"
        state["generated_images"] = [
            {
                "id": asset_id,
                "origin_type": "ai_generated",
                "provider": "OpenAI Image API",
                "model": "gpt-image-test",
                "prompt": "Ilustração educativa de teste.",
                "filename": "coeria_slide_2_ia.png",
                "media_type": "image/png",
                "data_base64": base64.b64encode(image_buffer.getvalue()).decode(
                    "ascii"
                ),
                "alt_text": "Ilustração educativa gerada por IA.",
                "approved": False,
                "created_at": "2026-08-18 00:00:00 UTC",
            }
        ]
        state["ai_image_generation_enabled"] = True
        slide = state["resources"]["presentation_outline"][1]
        slide["visual_mode"] = "ia"
        slide["visual_asset_id"] = asset_id
        slide["visual_prompt"] = "Representar visualmente o conceito sem texto."
        slide["visual_source"] = (
            "Imagem gerada por IA — OpenAI Image API, modelo gpt-image-test."
        )
        slide["alt_text"] = "Ilustração educativa gerada por IA."
        state["resources"]["quality"] = evaluate_quality(
            state, state["resources"]
        )
        self.assertTrue(state["resources"]["quality"]["passed"])

        state = review_current_stage(state, "approve", agent=self.agent)
        self.assertTrue(state["generated_images"][0]["approved"])
        state = review_current_stage(state, "approve", agent=self.agent)
        self.assertEqual(state["status"], "completed")

        presentation_path = Path(export_presentation(state))
        try:
            presentation = Presentation(presentation_path)
            pictures = [
                shape
                for shape in presentation.slides[1].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            self.assertEqual(len(pictures), 1)
            slide_text = "\n".join(
                shape.text
                for shape in presentation.slides[1].shapes
                if hasattr(shape, "text_frame")
            )
            self.assertIn("OpenAI Image API", slide_text)
            self.assertIn("gpt-image-test", slide_text)
        finally:
            presentation_path.unlink(missing_ok=True)

    def test_quality_rejects_unknown_document_image_id(self) -> None:
        state = self._resource_state()
        tampered = deepcopy(state)
        slide = tampered["resources"]["presentation_outline"][1]
        slide["visual_mode"] = "documento"
        slide["visual_asset_id"] = "document-inexistente"
        report = evaluate_quality(tampered, tampered["resources"])
        visual_check = next(
            item for item in report["checks"] if item["id"] == "presentation_visuals"
        )
        self.assertFalse(report["passed"])
        self.assertEqual(visual_check["status"], "error")

    def test_quality_validation_is_independent_from_the_generator(self) -> None:
        state = self._resource_state()
        tampered = deepcopy(state)
        tampered["teaching_activities"].pop()
        report = evaluate_quality(tampered, tampered["resources"])
        self.assertFalse(report["passed"])
        self.assertGreater(report["summary"]["errors"], 0)

    def test_quality_rejects_an_incomplete_visual_specification(self) -> None:
        state = self._resource_state()
        tampered = deepcopy(state)
        tampered["resources"]["presentation_outline"][1]["alt_text"] = ""
        report = evaluate_quality(tampered, tampered["resources"])
        visual_check = next(
            item for item in report["checks"] if item["id"] == "presentation_visuals"
        )
        self.assertFalse(report["passed"])
        self.assertEqual(visual_check["status"], "error")

    def test_quality_rejects_a_presentation_without_the_assessment_overview(self) -> None:
        state = self._resource_state()
        tampered = deepcopy(state)
        tampered["resources"]["presentation_outline"] = [
            slide
            for slide in tampered["resources"]["presentation_outline"]
            if not slide["title"].startswith(PRESENTATION_ASSESSMENT_TITLE)
        ]

        report = evaluate_quality(tampered, tampered["resources"])
        assessment_check = next(
            item
            for item in report["checks"]
            if item["id"] == "presentation_assessment_overview"
        )

        self.assertFalse(report["passed"])
        self.assertEqual(assessment_check["status"], "error")
        self.assertIn(PRESENTATION_ASSESSMENT_TITLE, assessment_check["detail"])

    def test_quality_explains_the_visual_failure_for_each_slide(self) -> None:
        state = self._resource_state()
        visual_items_schema = _schema_for("resources", state)["properties"][
            "artifact"
        ]["properties"]["presentation_outline"]["items"]["properties"][
            "visual_items"
        ]
        self.assertEqual(visual_items_schema["minItems"], 2)
        self.assertEqual(visual_items_schema["maxItems"], 4)
        tampered = deepcopy(state)
        slide = tampered["resources"]["presentation_outline"][0]
        slide["visual_items"] = ["A", "B", "C", "D", "E"]
        report = evaluate_quality(tampered, tampered["resources"])
        visual_check = next(
            item for item in report["checks"] if item["id"] == "presentation_visuals"
        )

        self.assertEqual(visual_check["status"], "error")
        self.assertIn("slide 1", visual_check["detail"])
        self.assertEqual(visual_check["target_stage"], "resources")
        self.assertEqual(visual_check["target_key"], "SLIDE:1")
        self.assertIn(
            "5 elementos; o diagrama admite 2 a 4",
            visual_check["detail"],
        )
        with self.assertRaisesRegex(
            AgentGenerationError,
            "5 elementos; o diagrama admite 2 a 4",
        ):
            validate_artifact("resources", tampered["resources"], tampered)

    def test_empty_artifacts_never_produce_misleading_green_checks(self) -> None:
        state = self._resource_state()
        tampered = deepcopy(state)
        tampered["learning_outcomes"] = []
        tampered["assessment_activities"] = []
        tampered["teaching_activities"] = []
        tampered["pedagogical_design"] = {"lessons": []}
        report = evaluate_quality(tampered, tampered["resources"])
        checks = {item["id"]: item for item in report["checks"]}
        affected = {
            "taxonomy_outcomes",
            "assessment_coverage",
            "assessment_purposes",
            "assessment_teaching_alignment",
            "teaching_coverage",
            "formative_activity_structure",
            "constructive_alignment",
        }
        self.assertTrue(all(checks[item]["status"] != "pass" for item in affected))
        self.assertEqual(checks["unique_outcomes"]["status"], "error")
        self.assertGreater(report["summary"]["warnings"], 0)

        final_validation = build_final_validation(tampered)
        sequence_check = next(
            item
            for item in final_validation["checks"]
            if item["id"] == "stage_pedagogical_design"
        )
        self.assertFalse(final_validation["passed"])
        self.assertFalse(sequence_check["passed"])
        self.assertEqual(sequence_check["target_stage"], "pedagogical_design")
        self.assertEqual(sequence_check["target_key"], "__stage__")
        self.assertIn("pelo menos uma aula", sequence_check["detail"])

    def test_missing_direct_assessment_link_navigates_to_the_task(self) -> None:
        state = self._resource_state()
        outcome_id = state["learning_outcomes"][0]["id"]
        for item in state["assessment_activities"]:
            item["outcome_ids"] = [
                identifier
                for identifier in item["outcome_ids"]
                if identifier != outcome_id
            ]

        report = evaluate_quality(state, state["resources"])
        coverage = next(
            item for item in report["checks"]
            if item["id"] == "assessment_coverage"
        )

        self.assertEqual(coverage["status"], "error")
        self.assertEqual(coverage["target_stage"], "assessment_activities")
        self.assertEqual(coverage["target_key"], outcome_id)

    def test_resources_with_blocking_quality_errors_cannot_be_approved(self) -> None:
        state = self._resource_state()
        state["resources"]["quality"]["passed"] = False
        with self.assertRaises(ValueError):
            review_current_stage(state, "approve", agent=self.agent)

    def test_resource_quality_failure_triggers_a_bounded_automatic_revision(self) -> None:
        alignment_state = create_session(
            self.course,
            resource_types=[RESOURCE_PRESENTATION, RESOURCE_TEST],
            agent=self.agent,
        )
        for _ in range(4):
            alignment_state = review_current_stage(
                alignment_state, "approve", agent=self.agent
            )

        class SelectiveRetryAgent:
            def __init__(self):
                self.delegate = create_test_agent()
                self.calls = []

            def generate(self, stage, state):
                resource_type = state["resource_types"][0]
                self.calls.append(resource_type)
                generation = self.delegate.generate(stage, state)
                artifact = deepcopy(generation.artifact)
                if (
                    resource_type == RESOURCE_PRESENTATION
                    and self.calls.count(resource_type) == 1
                ):
                    artifact["presentation_outline"][1]["outcome_id"] = (
                        "RA_INEXISTENTE"
                    )
                return GenerationResult(
                    artifact=artifact,
                    metadata={
                        "provider": "teste",
                        "model": "n/a",
                        "validation_attempts": 1,
                        "input_tokens": 10,
                        "output_tokens": 5,
                        "total_tokens": 15,
                    },
                )

        agent = SelectiveRetryAgent()
        progress = []
        with patch.dict(
            environ,
            {"COERIA_RESOURCE_QUALITY_MAX_REVISIONS": "1"},
            clear=False,
        ):
            result = review_current_stage(
                alignment_state,
                "approve",
                agent=agent,
                progress_callback=progress.append,
            )

        self.assertEqual(
            agent.calls,
            [RESOURCE_PRESENTATION, RESOURCE_PRESENTATION, RESOURCE_TEST],
        )
        self.assertTrue(result["resources"]["quality"]["passed"])
        resource_metadata = result["generation_metadata"]["resources"][-1][
            "resource_generations"
        ]
        self.assertEqual(len(resource_metadata), 2)
        self.assertEqual(resource_metadata[0]["quality_revisions"], 1)
        self.assertEqual(resource_metadata[1]["quality_revisions"], 0)
        self.assertEqual(len(resource_metadata[0]["attempts"]), 2)
        self.assertEqual(len(resource_metadata[1]["attempts"]), 1)
        self.assertEqual(
            result["generation_metadata"]["resources"][-1]["total_tokens"],
            45,
        )
        self.assertTrue(
            any(
                RESOURCE_PRESENTATION in item["event"]
                and "reformulado automaticamente" in item["event"]
                for item in result["audit"]
            )
        )
        self.assertTrue(
            any(
                f"A gerar recurso 1 de 2: {RESOURCE_PRESENTATION}" in message
                for message in progress
            )
        )
        self.assertTrue(
            any(
                f"A gerar recurso 2 de 2: {RESOURCE_TEST}" in message
                for message in progress
            )
        )

    def test_failed_resource_generation_resumes_from_persisted_drafts(self) -> None:
        alignment_state = create_session(
            self.course,
            resource_types=list(SUPPORTED_RESOURCE_TYPES),
            agent=self.agent,
        )
        for _ in range(4):
            alignment_state = review_current_stage(
                alignment_state,
                "approve",
                agent=self.agent,
            )

        class FailTestOnceAgent:
            def __init__(self):
                self.delegate = create_test_agent()
                self.calls = []
                self.failed = False

            def generate(self, stage, state):
                resource_type = state["resource_types"][0]
                self.calls.append(resource_type)
                if resource_type == RESOURCE_TEST and not self.failed:
                    self.failed = True
                    raise AgentGenerationError("Falha controlada no teste.")
                return self.delegate.generate(stage, state)

        agent = FailTestOnceAgent()
        first_progress = []
        second_progress = []
        with TemporaryDirectory() as temporary_directory:
            store = SQLiteSessionStore(
                Path(temporary_directory) / "resource-drafts.db"
            )
            service = ApplicationService(store)
            alignment_state["session_id"] = store.save(alignment_state)
            with patch(
                "prism.workflow.build_pedagogical_team",
                return_value=agent,
            ):
                with self.assertRaisesRegex(
                    AgentGenerationError,
                    "2 recursos já concluídos",
                ):
                    service.review_session(
                        alignment_state,
                        "approve",
                        progress_callback=first_progress.append,
                    )

                persisted = service.load_session(alignment_state["session_id"])
                self.assertEqual(persisted["current_stage"], "pedagogical_design")
                self.assertEqual(persisted["status"], "awaiting_review")
                self.assertNotIn("resources", persisted)
                self.assertEqual(
                    set(persisted["resource_generation_drafts"]["entries"]),
                    {RESOURCE_PRESENTATION, RESOURCE_WORKSHEET},
                )
                completed, _ = service.review_session(
                    persisted,
                    "approve",
                    progress_callback=second_progress.append,
                )

            stored_completed = service.load_session(completed["session_id"])

        self.assertEqual(
            agent.calls,
            [
                RESOURCE_PRESENTATION,
                RESOURCE_WORKSHEET,
                RESOURCE_TEST,
                RESOURCE_TEST,
                RESOURCE_PRACTICAL,
            ],
        )
        self.assertNotIn("resource_generation_drafts", completed)
        self.assertNotIn("resource_generation_drafts", stored_completed)
        self.assertTrue(completed["resources"]["quality"]["passed"])
        self.assertTrue(
            any(
                f"reutilizado: {RESOURCE_PRESENTATION}" in message
                for message in second_progress
            )
        )
        self.assertTrue(
            any(
                f"reutilizado: {RESOURCE_WORKSHEET}" in message
                for message in second_progress
            )
        )

    def test_completed_session_exports_a_complete_zip_package(self) -> None:
        state = review_current_stage(self._resource_state(), "approve", agent=self.agent)
        self.assertEqual(state["current_stage"], "final_validation")
        state = review_current_stage(state, "approve", agent=self.agent)
        package_path = Path(export_resource_package(state))
        try:
            with zipfile.ZipFile(package_path) as package:
                names = set(package.namelist())
                program_name = next(
                    name for name in names if name.endswith("_programa_uc.docx")
                )
                presentation_name = next(
                    name for name in names if name.endswith("_apresentacao.pptx")
                )
                self.assertTrue(any(name.endswith("_ficha_aula.docx") for name in names))
                self.assertTrue(any(name.endswith("_teste.docx") for name in names))
                self.assertTrue(any(name.endswith("_atividade_pratica.docx") for name in names))
                self.assertIn("sintese_alinhamento.csv", names)
                self.assertIn("rastreabilidade.csv", names)
                self.assertIn("manifesto.json", names)
                manifest = json.loads(package.read("manifesto.json"))
                self.assertEqual(manifest["application"], "CoerIA")
                self.assertIn("programas de unidades curriculares", manifest["application_name"])
                self.assertEqual(manifest["ai_provider"], "OpenAI")
                self.assertEqual(set(manifest["selected_resources"]), set(SUPPORTED_RESOURCE_TYPES))
                self.assertEqual(manifest["document_formats"], ["word"])
                self.assertEqual(manifest["primary_product"], program_name)
                self.assertEqual(manifest["primary_products"], [program_name])
                self.assertIn("visual_assets", manifest)
                self.assertIn("document_images", manifest["visual_assets"])
                self.assertIn("ai_generated_images", manifest["visual_assets"])
                program = Document(BytesIO(package.read(program_name)))
                program_text = "\n".join(
                    [paragraph.text for paragraph in program.paragraphs]
                    + [cell.text for table in program.tables for row in table.rows for cell in row.cells]
                )
                self.assertIn("Programa da Unidade Curricular", program_text)
                self.assertIn("Taxonomia selecionada", program_text)
                self.assertIn("Síntese automática do alinhamento", program_text)
                self.assertIn("Tarefas e critérios de avaliação", program_text)
                self.assertIn(
                    state["assessment_activities"][0]["teaching_activity_ids"][0],
                    program_text,
                )
                self.assertIn("Teaching for Quality Learning", program_text)
                self.assertNotIn("Learning outcomes", program_text)
                alignment_header = package.read("sintese_alinhamento.csv").decode(
                    "utf-8-sig"
                ).splitlines()[0]
                self.assertIn("Conteúdos", alignment_header)
                self.assertIn(
                    "Atividades de ensino-aprendizagem",
                    alignment_header,
                )
                self.assertNotIn("Atividades formativas", alignment_header)
                self.assertNotIn("Objetivos", alignment_header)
                self.assertNotIn("Recursos", alignment_header)
                presentation = Presentation(BytesIO(package.read(presentation_name)))
                self.assertGreaterEqual(len(presentation.slides), 3)
                for slide in presentation.slides:
                    self.assertGreaterEqual(len(slide.shapes), 5)
                    slide_text = "\n".join(
                        shape.text for shape in slide.shapes if hasattr(shape, "text_frame")
                    )
                    self.assertIn("Fonte visual:", slide_text)
                    descriptions = [
                        properties.get("descr", "")
                        for shape in slide.shapes
                        for properties in shape._element.xpath(".//p:cNvPr")
                    ]
                    self.assertTrue(any(description.strip() for description in descriptions))
        finally:
            package_path.unlink(missing_ok=True)

    def test_program_exports_manual_teaching_fields_and_lesson_planning(self) -> None:
        state = self._resource_state()
        activity = state["teaching_activities"][0]
        activity.pop("method", None)
        activity["learning_context"] = "Presencial"
        activity["practice"] = "Aplicação manual orientada."
        activity["support"] = "Acompanhamento manual do docente."
        lesson = state["pedagogical_design"]["lessons"][0]

        with TemporaryDirectory() as temporary_directory:
            word_path = Path(temporary_directory) / "programa.docx"
            latex_path = Path(temporary_directory) / "programa.tex"
            export_program_document(state, word_path)
            export_program_latex(state, latex_path)

            program = Document(word_path)
            word_text = "\n".join(
                [paragraph.text for paragraph in program.paragraphs]
                + [
                    cell.text
                    for table in program.tables
                    for row in table.rows
                    for cell in row.cells
                ]
            )
            self.assertIn("Contexto", word_text)
            self.assertIn("Prática", word_text)
            self.assertIn("Acompanhamento", word_text)
            self.assertIn("Aplicação manual orientada.", word_text)
            self.assertIn("Acompanhamento manual do docente.", word_text)
            self.assertIn("Tema", word_text)
            self.assertIn("Descrição do tema", word_text)
            self.assertIn("ISCED-F 2013", word_text)
            self.assertIn("0613", word_text)
            self.assertIn(
                "Desenvolvimento e análise de software e aplicações",
                word_text,
            )
            self.assertIn("7. Planeamento das aulas", word_text)
            self.assertIn(str(lesson["duration_minutes"]), word_text)
            self.assertIn(lesson["session_type"], word_text)
            self.assertIn(lesson["notes"], word_text)
            self.assertIn(lesson["component_ids"][0], word_text)
            self.assertIn(
                state["assessment_activities"][0]["teaching_activity_ids"][0],
                word_text,
            )
            assessment_headers = [
                cell.text for cell in program.tables[4].rows[0].cells
            ]
            self.assertIn("Atividades de ensino-aprendizagem", assessment_headers)
            self.assertIn("Resultados", assessment_headers)

            for table in program.tables[1:]:
                for row in table.rows[1:]:
                    row_properties = row._tr.get_or_add_trPr()
                    self.assertIsNotNone(
                        row_properties.find(
                            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}cantSplit"
                        )
                    )

            latex = latex_path.read_text(encoding="utf-8")
            self.assertIn(r"\textbf{Contexto}", latex)
            self.assertIn(r"\textbf{Prática}", latex)
            self.assertIn(r"\textbf{Acompanhamento}", latex)
            self.assertIn("Aplicação manual orientada.", latex)
            self.assertIn("Acompanhamento manual do docente.", latex)
            self.assertIn("Tema", latex)
            self.assertIn("Descrição do tema", latex)
            self.assertIn("ISCED-F 2013", latex)
            self.assertIn("0613", latex)
            self.assertIn(
                "Desenvolvimento e análise de software e aplicações",
                latex,
            )
            self.assertIn(r"\section{Planeamento das aulas}", latex)
            self.assertIn(str(lesson["duration_minutes"]), latex)
            self.assertIn(lesson["session_type"], latex)
            self.assertIn(lesson["notes"], latex)
            self.assertIn(lesson["component_ids"][0], latex)
            self.assertIn(
                state["assessment_activities"][0]["teaching_activity_ids"][0],
                latex,
            )

    def test_zip_package_respects_word_latex_or_both_document_formats(self) -> None:
        state = review_current_stage(self._resource_state(), "approve", agent=self.agent)
        state = review_current_stage(state, "approve", agent=self.agent)
        cases = (
            (("word",), True, False),
            (("latex",), False, True),
            (("word", "latex"), True, True),
        )
        for formats, expects_word, expects_latex in cases:
            with self.subTest(formats=formats):
                package_path = Path(export_resource_package(state, formats))
                try:
                    with zipfile.ZipFile(package_path) as package:
                        names = set(package.namelist())
                        word_names = {
                            name for name in names if name.endswith(".docx")
                        }
                        latex_names = {
                            name for name in names if name.endswith(".tex")
                        }
                        self.assertEqual(bool(word_names), expects_word)
                        self.assertEqual(bool(latex_names), expects_latex)
                        if expects_word:
                            self.assertEqual(len(word_names), 4)
                        if expects_latex:
                            self.assertEqual(len(latex_names), 4)
                            program_name = next(
                                name
                                for name in latex_names
                                if name.endswith("_programa_uc.tex")
                            )
                            program = package.read(program_name).decode("utf-8")
                            normalized_program = program.replace("\r\n", "\n")
                            self.assertIn(r"\begin{document}", program)
                            self.assertIn(
                                r"\section{Tarefas e critérios de avaliação}",
                                program,
                            )
                            self.assertIn(
                                r"\section{Síntese automática do alinhamento}",
                                program,
                            )
                            self.assertIn(
                                r"\clearpage"
                                "\n"
                                r"\section{Síntese automática do alinhamento}",
                                normalized_program,
                            )
                            self.assertIn(r"Biggs, J., \& Tang", program)
                            self.assertIn(r"\end{document}", program)
                            test_name = next(
                                name
                                for name in latex_names
                                if name.endswith("_teste.tex")
                            )
                            test_document = package.read(test_name).decode("utf-8")
                            self.assertIn(
                                r"\textbf{Tipo:} Resposta estruturada\par",
                                test_document,
                            )
                            self.assertIn(
                                r"\textbf{Resultado associado:}",
                                test_document,
                            )
                        self.assertTrue(
                            any(name.endswith("_apresentacao.pptx") for name in names)
                        )
                        manifest = json.loads(package.read("manifesto.json"))
                        self.assertEqual(manifest["document_formats"], list(formats))
                        self.assertEqual(
                            len(manifest["primary_products"]), len(formats)
                        )
                        self.assertEqual(
                            manifest["primary_product"],
                            manifest["primary_products"][0],
                        )
                finally:
                    package_path.unlink(missing_ok=True)

    def test_application_service_persists_export_format_choice_in_audit(self) -> None:
        state = review_current_stage(self._resource_state(), "approve", agent=self.agent)
        state = review_current_stage(state, "approve", agent=self.agent)
        state["resources"]["quality"]["passed"] = False
        with TemporaryDirectory() as temporary_directory:
            store = SQLiteSessionStore(
                Path(temporary_directory) / "coeria-export-formats.db"
            )
            service = ApplicationService(store)
            session_id = store.save(state)
            stored_state = service.load_session(session_id)
            package_data, package_filename, updated = service.export_session(
                stored_state,
                ("latex",),
            )
            self.assertTrue(package_filename.endswith(".zip"))
            with zipfile.ZipFile(BytesIO(package_data), "r") as package:
                self.assertIn("manifesto.json", package.namelist())
            self.assertEqual(updated["last_export_document_formats"], ["latex"])
            self.assertTrue(updated["resources"]["quality"]["passed"])
            self.assertEqual(
                updated["audit"][-1]["feedback"],
                "Formatos documentais: LaTeX.",
            )
            restored = service.load_session(session_id)
            self.assertEqual(
                restored["last_export_document_formats"],
                ["latex"],
            )
            self.assertTrue(restored["resources"]["quality"]["passed"])

    def test_latex_lists_use_a_teacher_fallback_when_empty(self) -> None:
        fallback = r"\emph{A confirmar pelo docente.}"
        self.assertEqual(_latex_itemize([]), fallback)
        self.assertEqual(_latex_itemize([], ordered=True), fallback)
        self.assertEqual(_latex_itemize(["", "  "]), fallback)

    def test_latex_list_items_escape_ambiguous_characters(self) -> None:
        document = _latex_itemize(['[Guia] <teste> "seguro"'])
        self.assertIn(
            r"\item{} {[}Guia{]} \textless{}teste\textgreater{} "
            r"\textquotedbl{}seguro\textquotedbl{}",
            document,
        )

    def test_latex_pdf_compiler_is_invoked_without_shell_escape(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            source = Path(temporary_directory) / "documento.tex"
            source.write_text(
                "\\documentclass{article}\\begin{document}Teste\\end{document}",
                encoding="utf-8",
            )

            def fake_run(command, **kwargs):
                source.with_suffix(".pdf").write_bytes(b"%PDF-1.7\n%%EOF\n")
                self.assertNotIn("shell", kwargs)
                return type("Result", (), {"returncode": 0, "stdout": b"OK"})()

            with (
                patch.dict(
                    environ,
                    {
                        "COERIA_LATEX_PDF_ENABLED": "true",
                        "COERIA_LATEX_COMPILER": "pdflatex",
                    },
                    clear=False,
                ),
                patch("prism.exporter.shutil.which", return_value="/usr/bin/pdflatex"),
                patch("prism.exporter.subprocess.run", side_effect=fake_run) as compiler,
            ):
                destination = Path(compile_latex_pdf(source) or "")

            self.assertEqual(compiler.call_count, 2)
            for call in compiler.call_args_list:
                command = call.args[0]
                self.assertIn("-no-shell-escape", command)
                self.assertIn("-halt-on-error", command)
                self.assertEqual(command[-1], source.name)
            self.assertTrue(destination.samefile(source.with_suffix(".pdf")))

    def test_latex_compiler_failure_is_written_to_the_server_log(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            source = Path(temporary_directory) / "documento.tex"
            source.write_text(
                "\\documentclass{article}\\begin{document}Teste\\end{document}",
                encoding="utf-8",
            )
            failed = type(
                "Result",
                (),
                {
                    "returncode": 1,
                    "stdout": b"! Undefined control sequence.\nline 12",
                },
            )()
            with (
                patch.dict(
                    environ,
                    {
                        "COERIA_LATEX_PDF_ENABLED": "true",
                        "COERIA_LATEX_COMPILER": "pdflatex",
                    },
                    clear=False,
                ),
                patch("prism.exporter.shutil.which", return_value="/usr/bin/pdflatex"),
                patch("prism.exporter.subprocess.run", return_value=failed),
                self.assertLogs("prism.exporter", level="ERROR") as captured,
            ):
                with self.assertRaisesRegex(ValueError, "Não foi possível compilar"):
                    compile_latex_pdf(source)

            log = "\n".join(captured.output)
            self.assertIn("passagem 1/2", log)
            self.assertIn("Undefined control sequence", log)

    def test_enabled_latex_pdf_compilation_requires_the_compiler(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            source = Path(temporary_directory) / "documento.tex"
            source.write_text(
                "\\documentclass{article}\\begin{document}Teste\\end{document}",
                encoding="utf-8",
            )
            with (
                patch.dict(
                    environ,
                    {"COERIA_LATEX_PDF_ENABLED": "true"},
                    clear=False,
                ),
                patch("prism.exporter.shutil.which", return_value=None),
            ):
                with self.assertRaisesRegex(ValueError, "não está disponível"):
                    compile_latex_pdf(source)

    def test_enabled_latex_pdf_compilation_adds_pdf_companions(self) -> None:
        state = review_current_stage(self._resource_state(), "approve", agent=self.agent)
        state = review_current_stage(state, "approve", agent=self.agent)

        def fake_compile(source_path):
            destination = Path(source_path).with_suffix(".pdf")
            destination.write_bytes(b"%PDF-1.7\n%%EOF\n")
            return str(destination)

        with (
            patch.dict(
                environ,
                {"COERIA_LATEX_PDF_ENABLED": "true"},
                clear=False,
            ),
            patch("prism.exporter._latex_compiler_path", return_value="/usr/bin/pdflatex"),
            patch("prism.exporter.compile_latex_pdf", side_effect=fake_compile),
        ):
            package_path = Path(export_resource_package(state, ("latex",)))
        try:
            with zipfile.ZipFile(package_path) as package:
                names = set(package.namelist())
                pdf_names = {name for name in names if name.endswith(".pdf")}
                self.assertEqual(len(pdf_names), 4)
                manifest = json.loads(package.read("manifesto.json"))
                self.assertTrue(manifest["latex_pdf_compilation"]["enabled"])
                self.assertEqual(
                    set(manifest["latex_pdf_compilation"]["generated_files"]),
                    pdf_names,
                )
        finally:
            package_path.unlink(missing_ok=True)

    @unittest.skipUnless(
        shutil.which("pdflatex"),
        "pdflatex não está instalado neste ambiente",
    )
    def test_generated_program_latex_compiles_with_installed_pdflatex(self) -> None:
        state = self._resource_state()
        with TemporaryDirectory() as temporary_directory:
            source = Path(temporary_directory) / "programa_uc.tex"
            export_program_latex(state, source)
            with patch.dict(
                environ,
                {"COERIA_LATEX_PDF_ENABLED": "true"},
                clear=False,
            ):
                destination = Path(compile_latex_pdf(source) or "")
            self.assertTrue(destination.is_file())
            self.assertTrue(destination.read_bytes().startswith(b"%PDF-"))


if __name__ == "__main__":
    unittest.main()
