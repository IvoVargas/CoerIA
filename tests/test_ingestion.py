import unittest
from os import environ
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from prism.ingestion import (
    SourceIngestionError,
    build_raw_source_text,
    build_source_text,
    extract_source_images,
    recover_direct_source_text,
)


class SourceIngestionTests(unittest.TestCase):
    @staticmethod
    def _minimal_pdf(text: str) -> bytes:
        content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
        objects = [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
            b"<< /Length " + str(len(content)).encode("ascii") + b" >>\nstream\n" + content + b"\nendstream",
            b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        ]
        result = bytearray(b"%PDF-1.4\n")
        offsets = [0]
        for index, body in enumerate(objects, start=1):
            offsets.append(len(result))
            result.extend(f"{index} 0 obj\n".encode("ascii") + body + b"\nendobj\n")
        xref_offset = len(result)
        result.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
        result.extend(b"0000000000 65535 f \n")
        for offset in offsets[1:]:
            result.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
        result.extend(
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode("ascii")
        )
        return bytes(result)

    def test_direct_text_and_text_files_are_combined_with_provenance(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            path = Path(temporary_directory) / "apoio.md"
            path.write_text("Conteúdo complementar sobre avaliação.", encoding="utf-8")
            result = build_source_text("Conteúdos programáticos principais.", [path])

        self.assertIn("[Texto introduzido pelo docente]", result)
        self.assertIn("[Ficheiro: apoio.md]", result)
        self.assertIn("Conteúdo complementar", result)

    def test_direct_text_is_recovered_without_internal_provenance_tags(self) -> None:
        combined = (
            "[Texto introduzido pelo docente]\nTexto original do docente.\n\n"
            "[Ficheiro: apoio.md]\nConteúdo extraído do documento."
        )

        self.assertEqual(
            recover_direct_source_text(combined),
            "Texto original do docente.",
        )
        self.assertEqual(
            recover_direct_source_text("[Ficheiro: apoio.md]\nApenas ficheiro."),
            "",
        )

    def test_unsupported_files_are_rejected(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            path = Path(temporary_directory) / "dados.csv"
            path.write_text("a,b", encoding="utf-8")
            with self.assertRaises(SourceIngestionError):
                build_source_text("", [path])

    def test_pdf_docx_and_pptx_text_is_extracted(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            directory = Path(temporary_directory)
            pdf_path = directory / "fonte.pdf"
            pdf_path.write_bytes(self._minimal_pdf("Conteudo PDF"))

            docx_path = directory / "fonte.docx"
            document = Document()
            document.add_paragraph("Conteúdo DOCX")
            document.save(docx_path)

            pptx_path = directory / "fonte.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Conteúdo PPTX"
            presentation.save(pptx_path)

            result = build_source_text("", [pdf_path, docx_path, pptx_path])

        self.assertIn("Conteudo PDF", result)
        self.assertIn("Conteúdo DOCX", result)
        self.assertIn("Conteúdo PPTX", result)

    def test_images_are_extracted_from_reference_documents_with_provenance(
        self,
    ) -> None:
        with TemporaryDirectory() as temporary_directory:
            directory = Path(temporary_directory)

            image_path = directory / "figura.png"
            Image.new(
                "RGB",
                (640, 400),
                (20, 120, 160),
            ).save(image_path)

            pdf_path = directory / "fonte.pdf"
            Image.open(image_path).convert("RGB").save(
                pdf_path,
                "PDF",
            )

            docx_path = directory / "fonte.docx"
            document = Document()
            document.add_paragraph("Documento com imagem.")
            document.add_picture(str(image_path))
            document.save(docx_path)

            pptx_path = directory / "fonte.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[6]
            )
            slide.shapes.add_picture(
                str(image_path),
                Inches(1),
                Inches(1),
            )
            presentation.save(pptx_path)

            assets = extract_source_images(
                [pdf_path, docx_path, pptx_path]
            )

        self.assertEqual(len(assets), 3)

        pdf_asset = next(
            item
            for item in assets
            if item["source_file"] == "fonte.pdf"
        )
        self.assertEqual(pdf_asset["source_location"], "Página 1")

        docx_asset = next(
            item
            for item in assets
            if item["source_file"] == "fonte.docx"
        )
        self.assertEqual(docx_asset["source_location"], "")

        pptx_asset = next(
            item
            for item in assets
            if item["source_file"] == "fonte.pptx"
        )
        self.assertEqual(pptx_asset["source_location"], "Slide 1")

        for asset in assets:
            self.assertEqual(asset["origin_type"], "document")
            self.assertTrue(asset["filename"])
            self.assertTrue(asset["data_base64"])
            self.assertTrue(asset["thumbnail_base64"])
            self.assertEqual(asset["image_mode"], "RGB")
            self.assertGreaterEqual(asset["width_px"], 240)
            self.assertGreaterEqual(asset["height_px"], 160)
            self.assertFalse(asset["approved"])
            self.assertEqual(asset["alt_text"], "")

    def test_pdf_catalogue_is_balanced_across_pages_after_full_scan(self) -> None:
        import io
        import pymupdf

        def png_bytes(colour):
            buffer = io.BytesIO()
            Image.new("RGB", (500, 320), colour).save(buffer, format="PNG")
            return buffer.getvalue()

        with TemporaryDirectory() as temporary_directory:
            pdf_path = Path(temporary_directory) / "catalogo.pdf"
            document = pymupdf.open()
            page1 = document.new_page(width=600, height=800)
            page1.insert_image(pymupdf.Rect(40, 60, 230, 210), stream=png_bytes((180, 20, 20)))
            page1.insert_image(pymupdf.Rect(350, 500, 550, 680), stream=png_bytes((20, 180, 20)))
            page2 = document.new_page(width=600, height=800)
            page2.insert_image(pymupdf.Rect(100, 180, 420, 430), stream=png_bytes((20, 20, 180)))
            document.save(pdf_path)
            document.close()

            with patch.dict(environ, {"COERIA_MAX_EXTRACTED_IMAGES": "2"}, clear=False):
                assets = extract_source_images([pdf_path])

        self.assertEqual(len(assets), 2)
        self.assertEqual(
            {asset["source_location"] for asset in assets},
            {"Página 1", "Página 2"},
        )

    def test_pdf_small_fragments_are_filtered_and_composite_is_rendered(self) -> None:
        import io
        import pymupdf

        def png_bytes(size, colour):
            buffer = io.BytesIO()
            Image.new("RGB", size, colour).save(buffer, format="PNG")
            return buffer.getvalue()

        with TemporaryDirectory() as temporary_directory:
            directory = Path(temporary_directory)

            filtered_path = directory / "filtrado.pdf"
            document = pymupdf.open()
            page = document.new_page(width=600, height=800)
            page.insert_image(
                pymupdf.Rect(20, 20, 50, 50),
                stream=png_bytes((40, 40), (10, 10, 10)),
            )
            page.insert_image(
                pymupdf.Rect(100, 120, 420, 360),
                stream=png_bytes((640, 400), (100, 140, 180)),
            )
            document.save(filtered_path)
            document.close()
            filtered = extract_source_images([filtered_path])

            composite_path = directory / "composta.pdf"
            document = pymupdf.open()
            page = document.new_page(width=600, height=800)
            page.insert_image(
                pymupdf.Rect(100, 160, 260, 360),
                stream=png_bytes((500, 400), (180, 100, 80)),
            )
            page.insert_image(
                pymupdf.Rect(270, 160, 430, 360),
                stream=png_bytes((500, 400), (80, 100, 180)),
            )
            document.save(composite_path)
            document.close()
            composite = extract_source_images([composite_path])

        self.assertEqual(len(filtered), 1)
        self.assertEqual(filtered[0]["candidate_kind"], "embedded")
        self.assertEqual(filtered[0]["image_mode"], "RGB")

        self.assertEqual(len(composite), 1)
        self.assertEqual(composite[0]["candidate_kind"], "composite_render")
        self.assertEqual(composite[0]["composite_object_count"], 2)
        self.assertEqual(composite[0]["image_mode"], "RGB")

    def test_raw_source_can_exceed_normal_context_limit_for_later_reduction(self) -> None:
        with patch.dict(
            environ,
            {
                "COERIA_MAX_SOURCE_CHARS": "20",
                "COERIA_MAX_RAW_SOURCE_CHARS": "200",
            },
            clear=False,
        ):
            result = build_raw_source_text(
                "Texto deliberadamente superior ao limite normal do contexto."
            )

        self.assertGreater(len(result), 20)

    def test_configured_source_limit_is_enforced(self) -> None:
        with patch.dict(environ, {"PRISM_MAX_SOURCE_CHARS": "20"}):
            with self.assertRaises(SourceIngestionError):
                build_source_text("Texto deliberadamente superior ao limite.")


if __name__ == "__main__":
    unittest.main()
